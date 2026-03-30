#!/usr/bin/env python3
import argparse
import base64
import hashlib
import io
import json
import logging
import os
import re
import secrets
import sqlite3
import time
import traceback
import webbrowser
from http.server import BaseHTTPRequestHandler, ThreadingHTTPServer
from logging.handlers import RotatingFileHandler
from pathlib import Path
from urllib.parse import parse_qs, unquote, urlparse

# ── Logging configuration ──
_LOG_DIR = Path(__file__).parent / "logs"
_LOG_DIR.mkdir(exist_ok=True)

logger = logging.getLogger("dev-tools")
logger.setLevel(logging.DEBUG)

# File handler (10MB, max 5 backups)
_fh = RotatingFileHandler(
    _LOG_DIR / "server.log", maxBytes=10 * 1024 * 1024, backupCount=5, encoding="utf-8"
)
_fh.setLevel(logging.DEBUG)
_fh.setFormatter(logging.Formatter("%(asctime)s [%(levelname)s] %(message)s", datefmt="%Y-%m-%d %H:%M:%S"))
logger.addHandler(_fh)

# Error-only file handler
_eh = RotatingFileHandler(
    _LOG_DIR / "error.log", maxBytes=10 * 1024 * 1024, backupCount=5, encoding="utf-8"
)
_eh.setLevel(logging.ERROR)
_eh.setFormatter(logging.Formatter("%(asctime)s [%(levelname)s] %(message)s\n%(exc_info)s", datefmt="%Y-%m-%d %H:%M:%S"))
logger.addHandler(_eh)

import subprocess, sys
import urllib.request

# ── CDN library definitions ──
_VENDOR_DIR = Path(__file__).parent / "static" / "vendor"

CDN_LIBS = [
    {"name": "jquery", "file": "jquery.min.js", "npm": "jquery", "version": "4.0.0",
     "url_tpl": "https://cdn.jsdelivr.net/npm/jquery@{v}/dist/jquery.min.js"},
    {"name": "lucide", "file": "lucide.min.js", "npm": "lucide", "version": "0.460.0",
     "url_tpl": "https://cdn.jsdelivr.net/npm/lucide@{v}/dist/umd/lucide.min.js"},
    {"name": "marked", "file": "marked.min.js", "npm": "marked", "version": "latest",
     "url_tpl": "https://cdn.jsdelivr.net/npm/marked@{v}/lib/marked.umd.js"},
    {"name": "highlight.js", "file": "highlight.min.js", "npm": "highlight.js", "version": "11.9.0",
     "url_tpl": "https://cdn.jsdelivr.net/gh/highlightjs/cdn-release@{v}/build/highlight.min.js"},
    {"name": "highlight.js-css", "file": "github.min.css", "npm": "highlight.js", "version": "11.9.0",
     "url_tpl": "https://cdn.jsdelivr.net/gh/highlightjs/cdn-release@{v}/build/styles/github.min.css"},
    {"name": "mermaid", "file": "mermaid.min.js", "npm": "mermaid", "version": "11",
     "url_tpl": "https://cdn.jsdelivr.net/npm/mermaid@{v}/dist/mermaid.min.js"},
    {"name": "jspdf", "file": "jspdf.umd.min.js", "npm": "jspdf", "version": "2.5.2",
     "url_tpl": "https://cdn.jsdelivr.net/npm/jspdf@{v}/dist/jspdf.umd.min.js"},
    {"name": "html2pdf", "file": "html2pdf.bundle.min.js", "npm": "html2pdf.js", "version": "0.10.2",
     "url_tpl": "https://cdn.jsdelivr.net/npm/html2pdf.js@{v}/dist/html2pdf.bundle.min.js"},
    {"name": "mammoth", "file": "mammoth.browser.min.js", "npm": "mammoth", "version": "1.8.0",
     "url_tpl": "https://cdn.jsdelivr.net/npm/mammoth@{v}/mammoth.browser.min.js"},
    {"name": "diff_match_patch", "file": "diff_match_patch.js", "npm": "diff-match-patch", "version": "20121119",
     "url_tpl": "https://cdnjs.cloudflare.com/ajax/libs/diff_match_patch/20121119/diff_match_patch.js"},
]

_VENDOR_META_FILE = _VENDOR_DIR / "_meta.json"

def _load_vendor_meta():
    """Local vendor library version metadata"""
    if _VENDOR_META_FILE.exists():
        try:
            return json.loads(_VENDOR_META_FILE.read_text("utf-8"))
        except Exception:
            pass
    return {}

def _save_vendor_meta(meta):
    _VENDOR_DIR.mkdir(parents=True, exist_ok=True)
    _VENDOR_META_FILE.write_text(json.dumps(meta, indent=2, ensure_ascii=False), "utf-8")

def _check_npm_latest(package_name):
    """Query latest version from npm registry"""
    try:
        url = f"https://registry.npmjs.org/{package_name}/latest"
        req = urllib.request.Request(url, headers={"User-Agent": "dev-tools/1.0", "Accept": "application/json"})
        with urllib.request.urlopen(req, timeout=10) as resp:
            data = json.loads(resp.read())
            return data.get("version")
    except Exception:
        return None

def _auto_install(package):
    """Auto-install package if not available"""
    cmds = [
        [sys.executable, "-m", "pip", "install", package],
        [sys.executable, "-m", "pip", "install", "--break-system-packages", package],
    ]
    for cmd in cmds:
        try:
            subprocess.check_call(cmd, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
            return True
        except subprocess.CalledProcessError:
            continue
    print(f"[Warning] Auto-install failed for {package}. Please install manually: pip install {package}")
    return False

try:
    from cryptography.fernet import Fernet
except ImportError:
    if _auto_install("cryptography"):
        from cryptography.fernet import Fernet
    else:
        Fernet = None

try:
    from openai import OpenAI
except ImportError:
    if _auto_install("openai"):
        from openai import OpenAI
    else:
        OpenAI = None

try:
    import anthropic as _anthropic_mod
except ImportError:
    if _auto_install("anthropic"):
        import anthropic as _anthropic_mod
    else:
        _anthropic_mod = None

try:
    from google import genai as google_genai
except ImportError:
    if _auto_install("google-genai"):
        from google import genai as google_genai
    else:
        google_genai = None

try:
    import xai_sdk
except ImportError:
    if _auto_install("xai-sdk"):
        import xai_sdk
    else:
        xai_sdk = None


# ── AI Provider abstraction ──
AI_PROVIDERS = {
    "openai":  {"label": "OpenAI",  "default_model": "gpt-4.1-nano"},
    "gemini":  {"label": "Gemini",  "default_model": "gemini-2.0-flash-lite"},
    "claude":  {"label": "Claude",  "default_model": "claude-haiku-4-5-20251001"},
    "grok":    {"label": "Grok",    "default_model": "grok-3-mini-fast"},
}


def _get_ai_api_keys():
    """Get all configured AI API keys from DB, fallback to env vars."""
    keys = {}
    try:
        conn = get_conn()
        row = conn.execute("SELECT value FROM dev_settings WHERE key='ai_api_keys'").fetchone()
        conn.close()
        if row:
            keys = json.loads(row[0])
    except Exception:
        pass
    # Fallback to env vars
    env_map = {
        "openai": "OPENAI_API_KEY",
        "gemini": "GEMINI_API_KEY",
        "claude": "ANTHROPIC_API_KEY",
        "grok": "GROK_API_KEY",
    }
    for provider, env_key in env_map.items():
        if not keys.get(provider):
            val = os.getenv(env_key, "").strip()
            if val:
                keys[provider] = val
    return keys


def _get_available_providers():
    """Return list of providers that have API keys configured."""
    keys = _get_ai_api_keys()
    result = []
    for pid, meta in AI_PROVIDERS.items():
        if keys.get(pid):
            result.append({"id": pid, "label": meta["label"], "default_model": meta["default_model"]})
    return result


def _ai_chat(provider, system_prompt, user_prompt, model=None):
    """Unified AI chat call across providers. Returns response text."""
    keys = _get_ai_api_keys()
    api_key = keys.get(provider, "")
    if not api_key:
        raise ValueError(f"API key not configured for {provider}")

    meta = AI_PROVIDERS.get(provider)
    if not meta:
        raise ValueError(f"Unknown provider: {provider}")
    model = model or meta["default_model"]

    if provider == "openai":
        if OpenAI is None:
            raise RuntimeError("openai package not installed")
        client = OpenAI(api_key=api_key)
        resp = client.responses.create(
            model=model,
            instructions=system_prompt,
            input=user_prompt,
        )
        return resp.output_text.strip()

    elif provider == "gemini":
        if google_genai is None:
            raise RuntimeError("google-genai package not installed")
        client = google_genai.Client(api_key=api_key)
        resp = client.models.generate_content(
            model=model,
            contents=user_prompt,
            config=google_genai.types.GenerateContentConfig(
                system_instruction=system_prompt,
            ),
        )
        return resp.text.strip()

    elif provider == "claude":
        if _anthropic_mod is None:
            raise RuntimeError("anthropic package not installed")
        client = _anthropic_mod.Anthropic(api_key=api_key)
        resp = client.messages.create(
            model=model,
            max_tokens=4096,
            system=system_prompt,
            messages=[{"role": "user", "content": user_prompt}],
        )
        return resp.content[0].text.strip()

    elif provider == "grok":
        if xai_sdk is None:
            raise RuntimeError("xai-sdk package not installed")
        client = xai_sdk.Client(api_key=api_key)
        conversation = client.chat.create(model=model)
        conversation.add_system(system_prompt)
        resp = conversation.add_user(user_prompt)
        return resp.text.strip()

    else:
        raise ValueError(f"Unsupported provider: {provider}")


# ── Developer mode: session store ──
_dev_sessions = {}  # {token: True}


def _hash_password(password: str, salt: bytes) -> str:
    return hashlib.pbkdf2_hmac("sha256", password.encode(), salt, 100000).hex()


def _derive_fernet_key(password: str, salt: bytes) -> bytes:
    dk = hashlib.pbkdf2_hmac("sha256", password.encode(), salt, 100000, dklen=32)
    return base64.urlsafe_b64encode(dk)


ROOT = Path(__file__).resolve().parent
STATIC_DIR = ROOT / "static"
DB_PATH = ROOT / "dev-tool.db"
DEFAULT_HOST = "127.0.0.1"
DEFAULT_PORT = 8080


def load_env_file(path):
    if not path.exists() or not path.is_file():
        return
    for line in path.read_text(encoding="utf-8").splitlines():
        raw = line.strip()
        if not raw or raw.startswith("#") or "=" not in raw:
            continue
        key, value = raw.split("=", 1)
        key = key.strip()
        value = value.strip().strip('"').strip("'")
        if key and key not in os.environ:
            os.environ[key] = value


def parse_bool(value, default=False):
    if value is None:
        return default
    return str(value).strip().lower() in {"1", "true", "yes", "y", "on"}


def get_conn():
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    return conn


def init_db():
    # Auto migration: mock_server.db → dev-tool.db
    conn = get_conn()
    conn.execute(
        """
        CREATE TABLE IF NOT EXISTS mocks (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            name TEXT NOT NULL,
            method TEXT NOT NULL,
            path TEXT NOT NULL,
            request_json TEXT,
            response_status INTEGER NOT NULL DEFAULT 200,
            response_headers TEXT,
            response_body TEXT,
            created_at TEXT NOT NULL DEFAULT CURRENT_TIMESTAMP,
            updated_at TEXT NOT NULL DEFAULT CURRENT_TIMESTAMP
        )
        """
    )
    conn.execute(
        """
        CREATE TABLE IF NOT EXISTS traffic_logs (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            matched_mock_id INTEGER,
            matched INTEGER NOT NULL DEFAULT 0,
            method TEXT NOT NULL,
            path TEXT NOT NULL,
            request_headers TEXT,
            request_body TEXT,
            request_json TEXT,
            response_status INTEGER NOT NULL,
            response_headers TEXT,
            response_body TEXT,
            created_at TEXT NOT NULL DEFAULT CURRENT_TIMESTAMP
        )
        """
    )
    conn.execute(
        """
        CREATE TABLE IF NOT EXISTS json_saves (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            name TEXT NOT NULL,
            data TEXT NOT NULL,
            created_at TEXT NOT NULL DEFAULT CURRENT_TIMESTAMP,
            updated_at TEXT NOT NULL DEFAULT CURRENT_TIMESTAMP
        )
        """
    )
    conn.execute(
        """
        CREATE TABLE IF NOT EXISTS csv_saves (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            name TEXT NOT NULL,
            data TEXT NOT NULL,
            delimiter TEXT NOT NULL DEFAULT ',',
            encoding TEXT NOT NULL DEFAULT 'utf-8',
            col_widths TEXT,
            created_at TEXT NOT NULL DEFAULT CURRENT_TIMESTAMP,
            updated_at TEXT NOT NULL DEFAULT CURRENT_TIMESTAMP
        )
        """
    )
    conn.execute(
        """
        CREATE TABLE IF NOT EXISTS md_saves (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            name TEXT NOT NULL,
            content TEXT NOT NULL,
            created_at TEXT NOT NULL DEFAULT CURRENT_TIMESTAMP,
            updated_at TEXT NOT NULL DEFAULT CURRENT_TIMESTAMP
        )
        """
    )
    conn.execute(
        """
        CREATE TABLE IF NOT EXISTS md_versions (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            save_id INTEGER NOT NULL,
            content TEXT NOT NULL,
            archived INTEGER NOT NULL DEFAULT 0,
            created_at TEXT NOT NULL DEFAULT CURRENT_TIMESTAMP,
            FOREIGN KEY (save_id) REFERENCES md_saves(id) ON DELETE CASCADE
        )
        """
    )
    conn.execute(
        """
        CREATE TABLE IF NOT EXISTS dev_settings (
            key TEXT PRIMARY KEY,
            value TEXT NOT NULL,
            encrypted INTEGER NOT NULL DEFAULT 0
        )
        """
    )
    # Add archived column if missing (existing DB migration)
    try:
        conn.execute("SELECT archived FROM md_versions LIMIT 1")
    except sqlite3.OperationalError:
        conn.execute("ALTER TABLE md_versions ADD COLUMN archived INTEGER NOT NULL DEFAULT 0")
    # Add version_num column if missing (existing DB migration)
    try:
        conn.execute("SELECT version_num FROM md_versions LIMIT 1")
    except sqlite3.OperationalError:
        conn.execute("ALTER TABLE md_versions ADD COLUMN version_num INTEGER NOT NULL DEFAULT 0")
        # Assign version_num to existing data (ordered by id per save_id)
        rows = conn.execute("SELECT id, save_id FROM md_versions ORDER BY save_id, id").fetchall()
        counters = {}
        for r in rows:
            sid = r["save_id"]
            counters[sid] = counters.get(sid, 0) + 1
            conn.execute("UPDATE md_versions SET version_num=? WHERE id=?", (counters[sid], r["id"]))
    # Add next_version column to md_saves (ensures monotonically increasing version numbers)
    try:
        conn.execute("SELECT next_version FROM md_saves LIMIT 1")
    except sqlite3.OperationalError:
        conn.execute("ALTER TABLE md_saves ADD COLUMN next_version INTEGER NOT NULL DEFAULT 1")
        # Initialize existing data: set to MAX(version_num) + 1 per save
        for row in conn.execute(
            "SELECT save_id, MAX(version_num) as mv FROM md_versions GROUP BY save_id"
        ).fetchall():
            conn.execute("UPDATE md_saves SET next_version=? WHERE id=?", (row["mv"] + 1, row["save_id"]))
    # Add comment column to md_versions
    try:
        conn.execute("SELECT comment FROM md_versions LIMIT 1")
    except sqlite3.OperationalError:
        conn.execute("ALTER TABLE md_versions ADD COLUMN comment TEXT NOT NULL DEFAULT ''")
    # ── git_commit_templates table ──
    conn.execute(
        """
        CREATE TABLE IF NOT EXISTS md_proofread_results (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            save_id INTEGER,
            items TEXT NOT NULL,
            created_at TEXT NOT NULL DEFAULT CURRENT_TIMESTAMP
        )
        """
    )
    conn.execute(
        """
        CREATE TABLE IF NOT EXISTS git_commit_templates (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            repo_path TEXT NOT NULL,
            name TEXT NOT NULL,
            template TEXT NOT NULL,
            created_at TEXT NOT NULL DEFAULT CURRENT_TIMESTAMP
        )
        """
    )
    conn.execute(
        """
        CREATE TABLE IF NOT EXISTS dataai_saves (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            prompt TEXT NOT NULL,
            format TEXT NOT NULL DEFAULT 'csv',
            count INTEGER NOT NULL DEFAULT 10,
            result TEXT NOT NULL,
            created_at TEXT NOT NULL DEFAULT CURRENT_TIMESTAMP
        )
        """
    )
    # ── i18n translations table ──
    conn.execute(
        """
        CREATE TABLE IF NOT EXISTS i18n (
            lang TEXT NOT NULL,
            key TEXT NOT NULL,
            value TEXT NOT NULL,
            PRIMARY KEY (lang, key)
        )
        """
    )
    # Seed i18n from JSON files (INSERT OR IGNORE — adds new keys without overwriting user edits)
    _seed_i18n(conn)
    conn.commit()
    conn.close()


def _seed_i18n(conn):
    """Load seed translations from static/lang/*.json into DB."""
    lang_dir = os.path.join(os.path.dirname(__file__), "static", "lang")
    for filename in os.listdir(lang_dir):
        if not filename.endswith(".json"):
            continue
        lang_code = filename[:-5]  # e.g. "ko", "en", "ja"
        filepath = os.path.join(lang_dir, filename)
        try:
            with open(filepath, "r", encoding="utf-8") as f:
                data = json.load(f)
            for key, value in data.items():
                conn.execute(
                    "INSERT OR IGNORE INTO i18n (lang, key, value) VALUES (?, ?, ?)",
                    (lang_code, key, value),
                )
            logger.info("[i18n] Seeded %d keys for lang '%s'", len(data), lang_code)
        except Exception as e:
            logger.error("[i18n] Failed to seed '%s': %s", filename, e)


# ── Custom Plugin System ──
_custom_plugins = {}  # id -> {manifest, path, routes_module}


def discover_custom_plugins():
    """Scan custom/ directory, load manifests, import routes, create DB tables."""
    custom_dir = ROOT / "custom"
    if not custom_dir.exists():
        custom_dir.mkdir(exist_ok=True)
        return

    for entry in sorted(custom_dir.iterdir()):
        if not entry.is_dir() or entry.name.startswith(".") or entry.name.startswith("_"):
            continue
        manifest_path = entry / "manifest.json"
        if not manifest_path.exists():
            logger.warning("[Plugin] %s has no manifest.json, skipping", entry.name)
            continue
        try:
            manifest = json.loads(manifest_path.read_text("utf-8"))
        except Exception as e:
            logger.error("[Plugin] Failed to read manifest for %s: %s", entry.name, e)
            continue

        plugin_id = manifest.get("id", entry.name)
        plugin = {"manifest": manifest, "path": entry, "routes_module": None}

        # DB table creation
        if "db_tables" in manifest:
            _init_plugin_tables(manifest["db_tables"], plugin_id)

        # Route loading
        if manifest.get("has_routes", False):
            routes_path = entry / "routes.py"
            if routes_path.exists():
                try:
                    import importlib.util
                    spec = importlib.util.spec_from_file_location(
                        f"custom_plugin_{plugin_id}", str(routes_path))
                    mod = importlib.util.module_from_spec(spec)
                    spec.loader.exec_module(mod)
                    plugin["routes_module"] = mod
                    logger.info("[Plugin] Loaded routes for '%s'", plugin_id)
                except Exception as e:
                    logger.error("[Plugin] Failed to load routes for '%s': %s", plugin_id, e)

        # i18n seeding
        lang_dir = entry / "lang"
        if lang_dir.exists():
            _seed_plugin_i18n(plugin_id, lang_dir)

        _custom_plugins[plugin_id] = plugin
        logger.info("[Plugin] Discovered '%s' v%s", plugin_id, manifest.get("version", "?"))


def _init_plugin_tables(tables_def, plugin_id):
    conn = get_conn()
    for table_name, table_def in tables_def.items():
        if not table_name.startswith("custom_"):
            logger.warning("[Plugin:%s] Table '%s' must start with 'custom_', skipping", plugin_id, table_name)
            continue
        cols = ", ".join(f"{col} {dtype}" for col, dtype in table_def["columns"].items())
        sql = f"CREATE TABLE IF NOT EXISTS {table_name} ({cols})"
        conn.execute(sql)
    conn.commit()
    conn.close()


def _seed_plugin_i18n(plugin_id, lang_dir):
    """Seed translations from custom plugin lang/ directory."""
    conn = get_conn()
    for filename in os.listdir(lang_dir):
        if not filename.endswith(".json"):
            continue
        lang_code = filename[:-5]
        try:
            with open(lang_dir / filename, "r", encoding="utf-8") as f:
                data = json.load(f)
            for key, value in data.items():
                full_key = f"custom.{plugin_id}.{key}"
                conn.execute(
                    "INSERT OR REPLACE INTO i18n (lang, key, value) VALUES (?, ?, ?)",
                    (lang_code, full_key, value))
        except Exception as e:
            logger.error("[Plugin i18n] Failed to seed %s/%s: %s", plugin_id, filename, e)
    conn.commit()
    conn.close()


def _get_plugin_enabled_state():
    """Get enabled/disabled state dict for all plugins."""
    conn = get_conn()
    row = conn.execute("SELECT value FROM dev_settings WHERE key='custom_plugins_enabled'").fetchone()
    conn.close()
    if row:
        try:
            return json.loads(row[0])
        except Exception:
            pass
    return {}


def _set_plugin_enabled_state(state):
    conn = get_conn()
    conn.execute(
        "INSERT OR REPLACE INTO dev_settings (key, value) VALUES (?, ?)",
        ("custom_plugins_enabled", json.dumps(state)))
    conn.commit()
    conn.close()


def parse_json_or_none(raw):
    if raw is None:
        return None
    if isinstance(raw, str):
        if raw.strip() == "":
            return None
        return json.loads(raw)
    return raw


def canonical_json(value):
    return json.dumps(value, ensure_ascii=False, sort_keys=True, separators=(",", ":"))


class MockHandler(BaseHTTPRequestHandler):
    server_version = "MockServer/1.0"

    def _read_body(self):
        length = int(self.headers.get("Content-Length", 0))
        if length == 0:
            return ""
        return self.rfile.read(length).decode("utf-8")

    def _read_body_bytes(self):
        length = int(self.headers.get("Content-Length", 0))
        if length == 0:
            return b""
        return self.rfile.read(length)

    def _parse_multipart(self):
        """Extract first file data from multipart/form-data"""
        content_type = self.headers.get("Content-Type", "")
        if "boundary=" not in content_type:
            return None
        boundary = content_type.split("boundary=")[1].strip()
        body = self._read_body_bytes()
        boundary_bytes = ("--" + boundary).encode()
        parts = body.split(boundary_bytes)
        for part in parts:
            if b"filename=" in part:
                # Separate header and body (delimited by blank line)
                header_end = part.find(b"\r\n\r\n")
                if header_end == -1:
                    continue
                file_data = part[header_end + 4:]
                # Strip trailing \r\n
                if file_data.endswith(b"\r\n"):
                    file_data = file_data[:-2]
                if file_data.endswith(b"--\r\n"):
                    file_data = file_data[:-4]
                if file_data.endswith(b"--"):
                    file_data = file_data[:-2]
                return file_data
        return None

    def _send_json(self, payload, status=200):
        if status >= 400:
            logger.warning(f"[API] {self.command} {self.path} → {status}: {payload}")
        body = json.dumps(payload, ensure_ascii=False).encode("utf-8")
        self.send_response(status)
        self.send_header("Content-Type", "application/json; charset=utf-8")
        self.send_header("Content-Length", str(len(body)))
        self.end_headers()
        self.wfile.write(body)

    def _send_text(self, text, status=200, content_type="text/plain; charset=utf-8"):
        data = text.encode("utf-8")
        self.send_response(status)
        self.send_header("Content-Type", content_type)
        self.send_header("Content-Length", str(len(data)))
        self.end_headers()
        self.wfile.write(data)

    def _serve_static(self, path):
        file_path = STATIC_DIR / path
        if path == "" or path == "/":
            file_path = STATIC_DIR / "index.html"
        elif path.startswith("/"):
            file_path = STATIC_DIR / path[1:]

        if not file_path.exists() or not file_path.is_file():
            if path.startswith("/vendor/"):
                fallback = STATIC_DIR / "vendor-default" / path[len("/vendor/"):]
                if fallback.exists() and fallback.is_file():
                    file_path = fallback
                else:
                    return False
            else:
                return False

        ext = file_path.suffix.lower()
        ctype = {
            ".html": "text/html; charset=utf-8",
            ".js": "application/javascript; charset=utf-8",
            ".css": "text/css; charset=utf-8",
            ".json": "application/json; charset=utf-8",
        }.get(ext, "application/octet-stream")
        data = file_path.read_bytes()
        self.send_response(200)
        self.send_header("Content-Type", ctype)
        self.send_header("Content-Length", str(len(data)))
        self.send_header("Cache-Control", "no-cache, no-store, must-revalidate")
        self.end_headers()
        self.wfile.write(data)
        return True

    def _list_mocks(self):
        conn = get_conn()
        rows = conn.execute(
            "SELECT id, name, method, path, request_json, response_status, response_headers, response_body, created_at, updated_at FROM mocks ORDER BY id DESC"
        ).fetchall()
        conn.close()

        items = []
        for r in rows:
            items.append(
                {
                    "id": r["id"],
                    "name": r["name"],
                    "method": r["method"],
                    "path": r["path"],
                    "request_json": json.loads(r["request_json"]) if r["request_json"] else None,
                    "response_status": r["response_status"],
                    "response_headers": json.loads(r["response_headers"]) if r["response_headers"] else {},
                    "response_body": json.loads(r["response_body"]) if r["response_body"] else None,
                    "created_at": r["created_at"],
                    "updated_at": r["updated_at"],
                }
            )
        self._send_json({"items": items})

    def _list_logs(self, limit=200):
        limit = max(1, min(int(limit), 1000))
        conn = get_conn()
        rows = conn.execute(
            """
            SELECT id, matched_mock_id, matched, method, path, request_headers, request_body, request_json,
                   response_status, response_headers, response_body, created_at
            FROM traffic_logs
            ORDER BY id DESC
            LIMIT ?
            """,
            (limit,),
        ).fetchall()
        conn.close()

        items = []
        for r in rows:
            items.append(
                {
                    "id": r["id"],
                    "matched_mock_id": r["matched_mock_id"],
                    "matched": bool(r["matched"]),
                    "method": r["method"],
                    "path": r["path"],
                    "request_headers": json.loads(r["request_headers"]) if r["request_headers"] else {},
                    "request_body": r["request_body"],
                    "request_json": json.loads(r["request_json"]) if r["request_json"] else None,
                    "response_status": r["response_status"],
                    "response_headers": json.loads(r["response_headers"]) if r["response_headers"] else {},
                    "response_body": json.loads(r["response_body"]) if r["response_body"] else None,
                    "created_at": r["created_at"],
                }
            )
        self._send_json({"items": items})

    def _clear_logs(self):
        conn = get_conn()
        conn.execute("DELETE FROM traffic_logs")
        conn.commit()
        conn.close()
        self._send_json({"ok": True})

    def _log_traffic(
        self,
        *,
        matched,
        matched_mock_id,
        method,
        path,
        request_headers,
        request_body,
        request_json,
        response_status,
        response_headers,
        response_body,
    ):
        conn = get_conn()
        conn.execute(
            """
            INSERT INTO traffic_logs (
                matched_mock_id, matched, method, path, request_headers, request_body, request_json,
                response_status, response_headers, response_body
            )
            VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
            """,
            (
                matched_mock_id,
                1 if matched else 0,
                method,
                path,
                canonical_json(request_headers or {}),
                request_body if request_body is not None else "",
                canonical_json(request_json) if request_json is not None else None,
                int(response_status),
                canonical_json(response_headers or {}),
                canonical_json(response_body) if response_body is not None else None,
            ),
        )
        conn.commit()
        conn.close()

    def _create_mock(self):
        raw = self._read_body()
        try:
            payload = json.loads(raw)
            name = payload["name"].strip()
            method = payload["method"].upper().strip()
            path = payload["path"].strip()
            if not path.startswith("/"):
                path = "/" + path

            request_json = parse_json_or_none(payload.get("request_json"))
            response_status = int(payload.get("response_status", 200))
            response_headers = payload.get("response_headers") or {}
            response_body = parse_json_or_none(payload.get("response_body"))

            if not name:
                raise ValueError("name is required")
            if not method:
                raise ValueError("method is required")
            if not path:
                raise ValueError("path is required")
            if not isinstance(response_headers, dict):
                raise ValueError("response_headers must be object")

            conn = get_conn()
            cur = conn.execute(
                """
                INSERT INTO mocks (name, method, path, request_json, response_status, response_headers, response_body)
                VALUES (?, ?, ?, ?, ?, ?, ?)
                """,
                (
                    name,
                    method,
                    path,
                    canonical_json(request_json) if request_json is not None else None,
                    response_status,
                    canonical_json(response_headers),
                    canonical_json(response_body) if response_body is not None else None,
                ),
            )
            conn.commit()
            new_id = cur.lastrowid
            conn.close()
            self._send_json({"ok": True, "id": new_id}, status=201)
        except (json.JSONDecodeError, KeyError, ValueError, TypeError) as e:
            self._send_json({"ok": False, "error": str(e)}, status=400)

    def _update_mock(self, mock_id):
        raw = self._read_body()
        try:
            payload = json.loads(raw)
            name = payload["name"].strip()
            method = payload["method"].upper().strip()
            path = payload["path"].strip()
            if not path.startswith("/"):
                path = "/" + path

            request_json = parse_json_or_none(payload.get("request_json"))
            response_status = int(payload.get("response_status", 200))
            response_headers = payload.get("response_headers") or {}
            response_body = parse_json_or_none(payload.get("response_body"))

            conn = get_conn()
            cur = conn.execute(
                """
                UPDATE mocks
                SET name=?, method=?, path=?, request_json=?, response_status=?, response_headers=?, response_body=?, updated_at=CURRENT_TIMESTAMP
                WHERE id=?
                """,
                (
                    name,
                    method,
                    path,
                    canonical_json(request_json) if request_json is not None else None,
                    response_status,
                    canonical_json(response_headers),
                    canonical_json(response_body) if response_body is not None else None,
                    mock_id,
                ),
            )
            conn.commit()
            changed = cur.rowcount
            conn.close()
            if changed == 0:
                self._send_json({"ok": False, "error": "not found"}, status=404)
                return
            self._send_json({"ok": True})
        except (json.JSONDecodeError, KeyError, ValueError, TypeError) as e:
            self._send_json({"ok": False, "error": str(e)}, status=400)

    def _delete_mock(self, mock_id):
        conn = get_conn()
        cur = conn.execute("DELETE FROM mocks WHERE id=?", (mock_id,))
        conn.commit()
        deleted = cur.rowcount
        conn.close()
        if deleted == 0:
            self._send_json({"ok": False, "error": "not found"}, status=404)
            return
        self._send_json({"ok": True})

    # --- Data AI ---
    def _dataai_generate(self):
        raw = self._read_body()
        try:
            payload = json.loads(raw)
            provider = payload.get("provider", "").strip()
            prompt = payload.get("prompt", "").strip()
            fmt = payload.get("format", "csv").strip()
            count = int(payload.get("count", 10))

            # Auto-select provider if not specified
            if not provider:
                available = _get_available_providers()
                if not available:
                    self._send_json({"ok": False, "error": "AI API 키가 설정되지 않았습니다. DEV > 일반 설정에서 API 키를 등록하세요."}, status=500)
                    return
                provider = available[0]["id"]

            if not prompt:
                self._send_json({"ok": False, "error": "데이터 설명을 입력하세요."}, status=400)
                return
            if count < 1 or count > 1000:
                self._send_json({"ok": False, "error": "건수는 1~1000 사이여야 합니다."}, status=400)
                return

            logger.info(f"[DataAI] Generation request: provider={provider}, prompt={prompt!r}, format={fmt}, count={count}")
            _t0 = time.time()

            system_prompt = "You are a data generator. Output ONLY a valid JSON array. No explanations, no markdown, no code blocks."

            user_prompt = f"""{prompt}

Generate EXACTLY {count} items as a JSON array of objects.
IMPORTANT rules:
- Each item must be UNIQUE and REALISTIC. No sequential/incrementing patterns.
- Do NOT repeat values with minor variations.
- Every value must be diverse and plausible.
- Generate values in the language specified or implied by the request above.
- Output ONLY the JSON array, nothing else."""

            result_text = _ai_chat(provider, system_prompt, user_prompt)
            elapsed = time.time() - _t0
            logger.info(f"[DataAI] AI response received: {elapsed:.1f}s, response length={len(result_text)}")
            all_items = self._parse_dataai_json(result_text)

            if not all_items:
                logger.error(f"[DataAI] JSON parsing failed. Raw response (first 500 chars): {result_text[:500]!r}")
                self._send_json({"ok": False, "error": "AI가 유효한 데이터를 생성하지 못했습니다."}, status=500)
                return

            logger.info(f"[DataAI] Parsing complete: requested={count}, actual={len(all_items)}")

            # Convert to requested format
            if fmt == "json":
                result = json.dumps(all_items, ensure_ascii=False, indent=2)
            elif fmt == "tsv":
                result = self._json_to_sv(all_items, "\t")
            else:  # csv
                result = self._json_to_sv(all_items, ",")

            # Auto-save to DB
            conn = get_conn()
            cur = conn.execute(
                "INSERT INTO dataai_saves (prompt, format, count, result) VALUES (?, ?, ?, ?)",
                (prompt, fmt, count, result),
            )
            conn.commit()
            save_id = cur.lastrowid
            conn.close()
            logger.info(f"[DataAI] Saved to DB: id={save_id}")

            self._send_json({"ok": True, "result": result, "count": len(all_items), "saved_id": save_id})
        except json.JSONDecodeError:
            logger.error("[DataAI] Request JSON parsing failed")
            self._send_json({"ok": False, "error": "잘못된 JSON 요청입니다."}, status=400)
        except BrokenPipeError:
            logger.warning("[DataAI] Client connection lost (BrokenPipe)")
        except Exception as e:
            logger.error(f"[DataAI] Generation failed: {e}", exc_info=True)
            try:
                self._send_json({"ok": False, "error": f"생성 실패: {e}"}, status=500)
            except BrokenPipeError:
                pass

    def _parse_dataai_json(self, text):
        """Parse JSON array from AI response"""
        text = text.strip()
        if text.startswith("```"):
            lines = text.split("\n")
            text = "\n".join(lines[1:-1] if lines[-1].strip() == "```" else lines[1:]).strip()
        try:
            data = json.loads(text)
            if isinstance(data, list):
                return data
        except json.JSONDecodeError:
            # Try extracting only the JSON array portion
            m = re.search(r'\[[\s\S]*\]', text)
            if m:
                try:
                    return json.loads(m.group())
                except json.JSONDecodeError:
                    pass
        return []

    def _json_to_sv(self, items, delimiter):
        """Convert JSON array to CSV/TSV string"""
        if not items:
            return ""
        keys = list(items[0].keys())
        lines = [delimiter.join(keys)]
        for item in items:
            row = []
            for k in keys:
                v = str(item.get(k, ""))
                if delimiter == "," and ("," in v or '"' in v or "\n" in v):
                    v = '"' + v.replace('"', '""') + '"'
                row.append(v)
            lines.append(delimiter.join(row))
        return "\n".join(lines)

    def _dataai_save(self):
        raw = self._read_body()
        try:
            payload = json.loads(raw)
            prompt = payload.get("prompt", "")
            fmt = payload.get("format", "csv")
            count = int(payload.get("count", 10))
            result = payload.get("result", "")
            if not result:
                self._send_json({"ok": False, "error": "저장할 데이터가 없습니다."}, status=400)
                return
            conn = get_conn()
            cur = conn.execute(
                "INSERT INTO dataai_saves (prompt, format, count, result) VALUES (?, ?, ?, ?)",
                (prompt, fmt, count, result),
            )
            conn.commit()
            save_id = cur.lastrowid
            conn.close()
            self._send_json({"ok": True, "id": save_id})
        except Exception as e:
            self._send_json({"ok": False, "error": str(e)}, status=500)

    def _dataai_list(self):
        conn = get_conn()
        rows = conn.execute("SELECT id, prompt, format, count, created_at FROM dataai_saves ORDER BY id DESC").fetchall()
        conn.close()
        self._send_json({"items": [dict(r) for r in rows]})

    def _dataai_get(self, save_id):
        conn = get_conn()
        row = conn.execute("SELECT * FROM dataai_saves WHERE id=?", (save_id,)).fetchone()
        conn.close()
        if not row:
            self._send_json({"error": "Not found"}, status=404)
            return
        self._send_json(dict(row))

    def _dataai_delete(self, save_id):
        conn = get_conn()
        conn.execute("DELETE FROM dataai_saves WHERE id=?", (save_id,))
        conn.commit()
        conn.close()
        self._send_json({"ok": True})

    LANG_NAMES = {"ko": "Korean", "en": "English", "ja": "Japanese"}

    def _translate(self):
        raw = self._read_body()
        try:
            payload = json.loads(raw)
            provider = payload.get("provider", "").strip()
            text = payload.get("text", "").strip()
            source = payload.get("source", "").strip()
            target = payload.get("target", "ko").strip()

            if not provider:
                available = _get_available_providers()
                if not available:
                    self._send_json({"ok": False, "error": "AI API 키가 설정되지 않았습니다. DEV > 일반 설정에서 API 키를 등록하세요."}, status=500)
                    return
                provider = available[0]["id"]

            if not text:
                self._send_json({"ok": False, "error": "번역할 텍스트를 입력하세요."}, status=400)
                return
            logger.info(f"[Translate] Request: provider={provider}, {source or 'auto'}→{target}, length={len(text)}")
            src_name = self.LANG_NAMES.get(source, "")
            tgt_name = self.LANG_NAMES.get(target, target)
            src_hint = f" The source language is {src_name}." if src_name else ""
            _t0 = time.time()
            result = _ai_chat(
                provider,
                f"You are a translator.{src_hint} Translate the user's text into {tgt_name}. Output ONLY the translated text, nothing else.",
                text,
            )
            logger.info(f"[Translate] Complete: {time.time()-_t0:.1f}s")
            self._send_json({"ok": True, "result": result})
        except json.JSONDecodeError:
            self._send_json({"ok": False, "error": "잘못된 JSON 요청입니다."}, status=400)
        except Exception as e:
            logger.error(f"[Translate] Failed: {e}", exc_info=True)
            self._send_json({"ok": False, "error": f"번역 실패: {e}"}, status=500)

    # --- Markdown AI proofreading ---
    def _md_proofread(self):
        raw = self._read_body()
        try:
            payload = json.loads(raw)
            provider = payload.get("provider", "").strip()
            text = payload.get("text", "").strip()
            include_style = payload.get("includeStyle", False)

            if not provider:
                available = _get_available_providers()
                if not available:
                    self._send_json({"ok": False, "error": "AI API 키가 설정되지 않았습니다. DEV > 일반 설정에서 API 키를 등록하세요."}, status=500)
                    return
                provider = available[0]["id"]

            if not text:
                self._send_json({"ok": False, "error": "검수할 텍스트를 입력하세요."}, status=400)
                return
            if len(text) > 50000:
                self._send_json({"ok": False, "error": f"문서가 너무 깁니다. ({len(text):,}자 / 제한 50,000자)"}, status=400)
                return

            logger.info(f"[Proofread] Request: provider={provider}, length={len(text)}, style_check={include_style}")
            _t0 = time.time()

            # Remove fenced code blocks
            code_block_pattern = re.compile(r"```[\s\S]*?```", re.MULTILINE)
            cleaned = code_block_pattern.sub("[CODE_BLOCK]", text)
            # Remove inline code
            inline_code_pattern = re.compile(r"`[^`]+`")
            cleaned = inline_code_pattern.sub("[CODE]", cleaned)

            # Generate text with line numbers
            lines = cleaned.split("\n")
            numbered = "\n".join(f"{i+1}: {line}" for i, line in enumerate(lines))

            style_instruction = ""
            if include_style:
                style_instruction = "\n- Style and readability issues (awkward phrasing, overly long sentences, unclear subjects)"

            prompt = f"""You are a proofreader for documents written in Korean, English, and Japanese.

The user will provide numbered lines of text extracted from a Markdown document.
Markdown syntax (headers, links, code blocks, etc.) should NOT be corrected.
[CODE_BLOCK] and [CODE] placeholders should be ignored completely.

Review ONLY the natural language text for:
- Spelling errors and typos{style_instruction}

Return a JSON array of corrections. Each item:
{{
  "line": <line number>,
  "before": "<original text fragment>",
  "after": "<corrected text fragment>",
  "reason": "<brief reason in the document's language>"
}}

If no corrections are needed, return an empty array: []
Return ONLY the JSON array, no other text."""

            result_text = _ai_chat(provider, prompt, numbered)
            logger.info(f"[Proofread] AI response: {time.time()-_t0:.1f}s")
            # Parse JSON
            try:
                items = json.loads(result_text)
                if not isinstance(items, list):
                    items = []
            except json.JSONDecodeError:
                logger.error(f"[Proofread] AI response parsing failed. Raw (first 500 chars): {result_text[:500]!r}")
                self._send_json({"ok": False, "error": "AI 응답 파싱 실패. 다시 시도해주세요."}, status=500)
                return
            logger.info(f"[Proofread] Complete: {len(items)} corrections")
            self._send_json({"ok": True, "items": items})
        except json.JSONDecodeError:
            self._send_json({"ok": False, "error": "잘못된 JSON 요청입니다."}, status=400)
        except Exception as e:
            logger.error(f"[Proofread] Failed: {e}", exc_info=True)
            self._send_json({"ok": False, "error": f"검수 실패: {e}"}, status=500)

    # --- Proofread result save/get/delete ---
    def _save_proofread_result(self):
        raw = self._read_body()
        try:
            payload = json.loads(raw)
            save_id = payload.get("save_id")  # nullable
            items = payload.get("items", [])
            conn = get_conn()
            # Delete existing result for the same save_id (keep only 1)
            if save_id is not None:
                conn.execute("DELETE FROM md_proofread_results WHERE save_id=?", (save_id,))
            cur = conn.execute(
                "INSERT INTO md_proofread_results (save_id, items) VALUES (?, ?)",
                (save_id, json.dumps(items, ensure_ascii=False)),
            )
            conn.commit()
            self._send_json({"ok": True, "id": cur.lastrowid})
        except Exception as e:
            self._send_json({"ok": False, "error": str(e)}, status=500)

    def _get_proofread_result(self, save_id):
        conn = get_conn()
        row = conn.execute(
            "SELECT id, save_id, items, created_at FROM md_proofread_results WHERE save_id=? ORDER BY id DESC LIMIT 1",
            (save_id,),
        ).fetchone()
        if not row:
            self._send_json({"ok": True, "result": None})
            return
        self._send_json({
            "ok": True,
            "result": {
                "id": row["id"],
                "save_id": row["save_id"],
                "items": json.loads(row["items"]),
                "created_at": row["created_at"],
            },
        })

    def _delete_proofread_result(self, result_id):
        conn = get_conn()
        conn.execute("DELETE FROM md_proofread_results WHERE id=?", (result_id,))
        conn.commit()
        self._send_json({"ok": True})

    # --- CSV saves CRUD ---
    def _list_csv_saves(self):
        conn = get_conn()
        rows = conn.execute(
            "SELECT id, name, delimiter, encoding, created_at, updated_at FROM csv_saves ORDER BY updated_at DESC"
        ).fetchall()
        conn.close()
        items = [
            {
                "id": r["id"],
                "name": r["name"],
                "delimiter": r["delimiter"],
                "encoding": r["encoding"],
                "created_at": r["created_at"],
                "updated_at": r["updated_at"],
            }
            for r in rows
        ]
        self._send_json({"items": items})

    def _get_csv_save(self, save_id):
        conn = get_conn()
        row = conn.execute(
            "SELECT * FROM csv_saves WHERE id=?", (save_id,)
        ).fetchone()
        conn.close()
        if not row:
            self._send_json({"ok": False, "error": "not found"}, status=404)
            return
        self._send_json({
            "id": row["id"],
            "name": row["name"],
            "data": json.loads(row["data"]),
            "delimiter": row["delimiter"],
            "encoding": row["encoding"],
            "col_widths": json.loads(row["col_widths"]) if row["col_widths"] else None,
            "created_at": row["created_at"],
            "updated_at": row["updated_at"],
        })

    def _create_csv_save(self):
        raw = self._read_body()
        try:
            payload = json.loads(raw)
            name = payload.get("name", "").strip()
            data = payload.get("data")
            delimiter = payload.get("delimiter", ",")
            encoding = payload.get("encoding", "utf-8")
            col_widths = payload.get("col_widths")
            if not name:
                self._send_json({"ok": False, "error": "이름을 입력하세요."}, status=400)
                return
            if not data or not isinstance(data, list):
                self._send_json({"ok": False, "error": "데이터가 없습니다."}, status=400)
                return
            conn = get_conn()
            cur = conn.execute(
                """
                INSERT INTO csv_saves (name, data, delimiter, encoding, col_widths)
                VALUES (?, ?, ?, ?, ?)
                """,
                (
                    name,
                    json.dumps(data, ensure_ascii=False),
                    delimiter,
                    encoding,
                    json.dumps(col_widths) if col_widths else None,
                ),
            )
            conn.commit()
            new_id = cur.lastrowid
            conn.close()
            self._send_json({"ok": True, "id": new_id}, status=201)
        except (json.JSONDecodeError, KeyError, ValueError) as e:
            self._send_json({"ok": False, "error": str(e)}, status=400)

    def _update_csv_save(self, save_id):
        raw = self._read_body()
        try:
            payload = json.loads(raw)
            name = payload.get("name", "").strip()
            data = payload.get("data")
            delimiter = payload.get("delimiter", ",")
            encoding = payload.get("encoding", "utf-8")
            col_widths = payload.get("col_widths")
            if not name:
                self._send_json({"ok": False, "error": "이름을 입력하세요."}, status=400)
                return
            conn = get_conn()
            cur = conn.execute(
                """
                UPDATE csv_saves
                SET name=?, data=?, delimiter=?, encoding=?, col_widths=?, updated_at=CURRENT_TIMESTAMP
                WHERE id=?
                """,
                (
                    name,
                    json.dumps(data, ensure_ascii=False),
                    delimiter,
                    encoding,
                    json.dumps(col_widths) if col_widths else None,
                    save_id,
                ),
            )
            conn.commit()
            changed = cur.rowcount
            conn.close()
            if changed == 0:
                self._send_json({"ok": False, "error": "not found"}, status=404)
                return
            self._send_json({"ok": True})
        except (json.JSONDecodeError, KeyError, ValueError) as e:
            self._send_json({"ok": False, "error": str(e)}, status=400)

    def _delete_csv_save(self, save_id):
        conn = get_conn()
        cur = conn.execute("DELETE FROM csv_saves WHERE id=?", (save_id,))
        conn.commit()
        deleted = cur.rowcount
        conn.close()
        if deleted == 0:
            self._send_json({"ok": False, "error": "not found"}, status=404)
            return
        self._send_json({"ok": True})

    # ── JSON saves ────────────────────────────────────────────

    def _list_json_saves(self):
        conn = get_conn()
        rows = conn.execute(
            "SELECT id, name, created_at, updated_at FROM json_saves ORDER BY updated_at DESC"
        ).fetchall()
        conn.close()
        items = [
            {"id": r["id"], "name": r["name"], "created_at": r["created_at"], "updated_at": r["updated_at"]}
            for r in rows
        ]
        self._send_json({"items": items})

    def _get_json_save(self, save_id):
        conn = get_conn()
        row = conn.execute("SELECT * FROM json_saves WHERE id=?", (save_id,)).fetchone()
        conn.close()
        if not row:
            self._send_json({"ok": False, "error": "not found"}, status=404)
            return
        self._send_json({
            "id": row["id"], "name": row["name"], "data": row["data"],
            "created_at": row["created_at"], "updated_at": row["updated_at"],
        })

    def _create_json_save(self):
        raw = self._read_body()
        try:
            payload = json.loads(raw)
            name = payload.get("name", "").strip()
            data = payload.get("data", "")
            if not name:
                self._send_json({"ok": False, "error": "이름을 입력하세요."}, status=400)
                return
            if not data:
                self._send_json({"ok": False, "error": "저장할 내용이 없습니다."}, status=400)
                return
            conn = get_conn()
            cur = conn.execute(
                "INSERT INTO json_saves (name, data) VALUES (?, ?)", (name, data)
            )
            conn.commit()
            new_id = cur.lastrowid
            conn.close()
            self._send_json({"ok": True, "id": new_id}, status=201)
        except (json.JSONDecodeError, KeyError, ValueError) as e:
            self._send_json({"ok": False, "error": str(e)}, status=400)

    def _update_json_save(self, save_id):
        raw = self._read_body()
        try:
            payload = json.loads(raw)
            name = payload.get("name", "").strip()
            data = payload.get("data", "")
            if not name:
                self._send_json({"ok": False, "error": "이름을 입력하세요."}, status=400)
                return
            conn = get_conn()
            cur = conn.execute(
                "UPDATE json_saves SET name=?, data=?, updated_at=CURRENT_TIMESTAMP WHERE id=?",
                (name, data, save_id),
            )
            conn.commit()
            changed = cur.rowcount
            conn.close()
            if changed == 0:
                self._send_json({"ok": False, "error": "not found"}, status=404)
                return
            self._send_json({"ok": True})
        except (json.JSONDecodeError, KeyError, ValueError) as e:
            self._send_json({"ok": False, "error": str(e)}, status=400)

    def _delete_json_save(self, save_id):
        conn = get_conn()
        cur = conn.execute("DELETE FROM json_saves WHERE id=?", (save_id,))
        conn.commit()
        deleted = cur.rowcount
        conn.close()
        if deleted == 0:
            self._send_json({"ok": False, "error": "not found"}, status=404)
            return
        self._send_json({"ok": True})

    # ── Markdown saves ──────────────────────────────────────────

    def _list_md_saves(self):
        conn = get_conn()
        rows = conn.execute(
            "SELECT id, name, created_at, updated_at FROM md_saves ORDER BY updated_at DESC"
        ).fetchall()
        conn.close()
        self._send_json([dict(r) for r in rows])

    def _get_md_save(self, save_id):
        conn = get_conn()
        row = conn.execute(
            "SELECT * FROM md_saves WHERE id=?", (save_id,)
        ).fetchone()
        conn.close()
        if not row:
            self._send_json({"ok": False, "error": "not found"}, status=404)
            return
        self._send_json(dict(row))

    def _get_md_save_html(self, save_id):
        """Return saved markdown as rendered HTML page (for popup view, client-side rendering via marked.js)"""
        conn = get_conn()
        row = conn.execute(
            "SELECT * FROM md_saves WHERE id=?", (save_id,)
        ).fetchone()
        conn.close()
        if not row:
            self.send_response(404)
            self.send_header("Content-Type", "text/html; charset=utf-8")
            self.end_headers()
            self.wfile.write(b"<h1>Not Found</h1>")
            return

        import html as html_mod
        title = html_mod.escape(row["name"])
        # JSON encode for safe insertion into JS string
        content_json = json.dumps(row["content"] or "")

        page = f"""<!DOCTYPE html>
<html><head>
<meta charset="utf-8">
<title>{title}</title>
<link rel="stylesheet" href="/vendor/github.min.css">
<style>
:root {{
  --accent: #3b82f6;
  --line: #e2e8f0;
  --bg: #f8fafc;
  --text-secondary: #64748b;
}}
body {{
  font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, sans-serif;
  max-width: 760px; margin: 0 auto; padding: 32px 24px;
  color: #1e293b; line-height: 1.7;
}}
/* Same styles as editor preview (.md-preview-pane) */
h1 {{ font-size: 1.8em; margin: 0.6em 0 0.4em; padding-bottom: 0.3em; border-bottom: 1px solid var(--line); }}
h2 {{ font-size: 1.5em; margin: 0.5em 0 0.3em; padding-bottom: 0.2em; border-bottom: 1px solid var(--line); }}
h3 {{ font-size: 1.25em; margin: 0.5em 0 0.3em; }}
h4, h5, h6 {{ font-size: 1.1em; margin: 0.5em 0 0.3em; }}
p {{ margin: 0.5em 0; }}
code {{ background: #f0f4f8; padding: 2px 6px; border-radius: 4px; font-size: 0.9em; font-family: 'SF Mono', Consolas, monospace; }}
pre {{ background: #f6f8fa; padding: 12px 16px; border-radius: 8px; overflow-x: auto; margin: 0.8em 0; box-sizing: border-box; }}
pre code {{ background: none; padding: 0; font-size: 13px; display: block; white-space: pre; }}
blockquote {{ border-left: 4px solid var(--accent); margin: 0.5em 0; padding: 8px 16px; color: var(--text-secondary); background: rgba(59,130,246,0.04); border-radius: 0 4px 4px 0; }}
table {{ width: 100%; border-collapse: collapse; margin: 0.8em 0; }}
th, td {{ border: 1px solid var(--line); padding: 6px 12px; text-align: left; }}
th {{ background: var(--bg); font-weight: 600; }}
img {{ max-width: 100%; height: auto; border-radius: 4px; }}
a {{ color: var(--accent); text-decoration: none; }}
a:hover {{ text-decoration: underline; }}
.anchor {{ color: #c0c8d0; font-size: 0.8em; margin-right: 4px; text-decoration: none; }}
.anchor:hover {{ color: var(--accent); }}
h1:target, h2:target, h3:target, h4:target, h5:target, h6:target {{ scroll-margin-top: 16px; }}
ul, ol {{ padding-left: 24px; margin: 0.5em 0; }}
li {{ margin: 0.2em 0; }}
hr {{ border: none; border-top: 3px solid #d0d7de; margin: 1.5em 0; }}
input[type="checkbox"] {{ margin-right: 6px; }}
/* frontmatter */
.md-frontmatter {{ margin-bottom: 16px; border: 1px solid #d8dee4; border-radius: 6px; overflow: hidden; background: #f6f8fa; }}
.md-frontmatter table {{ width: 100%; border-collapse: collapse; font-size: 13px; }}
.md-frontmatter th {{ text-align: left; padding: 6px 12px; background: #eaeef2; color: #24292f; font-weight: 600; width: 140px; border-bottom: 1px solid #d8dee4; }}
.md-frontmatter td {{ padding: 6px 12px; color: #57606a; border-bottom: 1px solid #d8dee4; }}
.md-frontmatter tr:last-child th, .md-frontmatter tr:last-child td {{ border-bottom: none; }}
/* GitHub-style alerts */
.md-alert {{ padding: 12px 16px; margin: 16px 0; border-left: 4px solid; border-radius: 6px; }}
.md-alert p {{ margin: 4px 0; }}
.md-alert-title {{ font-weight: 600; margin-bottom: 4px !important; }}
.md-alert-note {{ border-color: #0969da; background: #ddf4ff; }}
.md-alert-note .md-alert-title {{ color: #0969da; }}
.md-alert-tip {{ border-color: #1a7f37; background: #dafbe1; }}
.md-alert-tip .md-alert-title {{ color: #1a7f37; }}
.md-alert-important {{ border-color: #8250df; background: #fbefff; }}
.md-alert-important .md-alert-title {{ color: #8250df; }}
.md-alert-warning {{ border-color: #9a6700; background: #fff8c5; }}
.md-alert-warning .md-alert-title {{ color: #9a6700; }}
.md-alert-caution {{ border-color: #cf222e; background: #ffebe9; }}
.md-alert-caution .md-alert-title {{ color: #cf222e; }}
</style>
</head><body>
<div id="content"></div>
<script src="/vendor/marked.min.js"></script>
<script src="/vendor/highlight.min.js"></script>
<script src="/vendor/mermaid.min.js"></script>
<script>
(function() {{
  var raw = {content_json};

  // frontmatter parsing
  function parseFrontmatter(src) {{
    var match = src.match(/^---\\r?\\n([\\s\\S]*?)\\r?\\n---\\r?\\n?([\\s\\S]*)$/);
    if (!match) return {{ meta: null, body: src }};
    var lines = match[1].split("\\n"), meta = {{}};
    for (var i = 0; i < lines.length; i++) {{
      var idx = lines[i].indexOf(":");
      if (idx > 0) {{
        meta[lines[i].slice(0, idx).trim()] = lines[i].slice(idx + 1).trim();
      }}
    }}
    return {{ meta: meta, body: match[2] }};
  }}

  function renderFrontmatterHtml(meta) {{
    if (!meta || Object.keys(meta).length === 0) return "";
    var html = '<div class="md-frontmatter"><table>';
    for (var k in meta) {{
      html += "<tr><th>" + k.replace(/</g, "&lt;") + "</th><td>" + meta[k].replace(/</g, "&lt;") + "</td></tr>";
    }}
    return html + "</table></div>";
  }}

  // GitHub-style alerts post-processing
  var alertLabels = {{ NOTE: "참고", TIP: "팁", IMPORTANT: "중요", WARNING: "경고", CAUTION: "주의" }};
  var alertIcons = {{ NOTE: "ℹ️", TIP: "💡", IMPORTANT: "❗", WARNING: "⚠️", CAUTION: "🔴" }};
  function processAlerts(html) {{
    return html.replace(
      /<blockquote>\\s*<p>\\[!(NOTE|TIP|IMPORTANT|WARNING|CAUTION)\\](?:<br\\s*\\/?>|\\n)?\\s*([\\s\\S]*?)<\\/blockquote>/gi,
      function(match, type, rest) {{
        var t = type.toUpperCase();
        return '<div class="md-alert md-alert-' + t.toLowerCase() + '">'
          + '<p class="md-alert-title">' + alertIcons[t] + " " + alertLabels[t] + "</p>"
          + "<p>" + rest + "</div>";
      }}
    );
  }}

  // marked custom renderer (hljs code highlighting + mermaid)
  var mermaidId = 0;
  var renderer = new marked.Renderer();
  var origCode = renderer.code.bind(renderer);
  renderer.heading = function(token) {{
    var text = token.text || token;
    var depth = token.depth || 1;
    var slug = text.toLowerCase().replace(/<[^>]*>/g, "").replace(/[^\\w가-힣ㄱ-ㅎㅏ-ㅣ\\s-]/g, "").replace(/\\s+/g, "-").replace(/-+$/,"");
    return '<h' + depth + ' id="' + slug + '"><a class="anchor" href="#' + slug + '">#</a> ' + text + '</h' + depth + '>';
  }};
  renderer.code = function(code, lang, escaped) {{
    var text = code, language = lang;
    if (typeof code === "object" && code !== null) {{
      text = code.text || "";
      language = code.lang || lang;
    }}
    if (language === "mermaid") {{
      mermaidId++;
      return '<div class="mermaid" id="peek-mermaid-' + mermaidId + '">' + text + '</div>';
    }}
    if (typeof hljs !== "undefined") {{
      var highlighted;
      if (language && hljs.getLanguage(language)) {{
        highlighted = hljs.highlight(text, {{ language: language }}).value;
      }} else {{
        highlighted = hljs.highlightAuto(text).value;
      }}
      return '<pre><code class="hljs language-' + (language || "") + '">' + highlighted + '</code></pre>';
    }}
    return origCode(code, lang, escaped);
  }};

  marked.setOptions({{ renderer: renderer, breaks: true, gfm: true }});

  // rendering
  var parsed = parseFrontmatter(raw);
  var fmHtml = renderFrontmatterHtml(parsed.meta);
  document.getElementById("content").innerHTML = fmHtml + processAlerts(marked.parse(parsed.body));

  // mermaid rendering
  mermaid.initialize({{ startOnLoad: false, theme: "default" }});
  document.querySelectorAll(".mermaid").forEach(function(el) {{
    var code = el.textContent;
    el.dataset.processed = "true";
    mermaid.render("peek-svg-" + Date.now() + Math.random().toString(36).slice(2), code).then(function(result) {{
      el.innerHTML = result.svg;
    }}).catch(function(e) {{
      el.innerHTML = '<pre style="color:#bf233a;font-size:12px">Mermaid 오류: ' + (e.message || e) + '</pre>';
    }});
  }});
}})();
</script>
</body></html>"""

        body = page.encode("utf-8")
        self.send_response(200)
        self.send_header("Content-Type", "text/html; charset=utf-8")
        self.send_header("Content-Length", str(len(body)))
        self.end_headers()
        self.wfile.write(body)

    def _create_md_save(self):
        raw = self._read_body()
        try:
            payload = json.loads(raw)
            name = payload.get("name", "").strip()
            content = payload.get("content", "")
            if not name:
                self._send_json({"ok": False, "error": "이름을 입력하세요."}, status=400)
                return
            conn = get_conn()
            cur = conn.execute(
                "INSERT INTO md_saves (name, content) VALUES (?, ?)",
                (name, content),
            )
            conn.commit()
            new_id = cur.lastrowid
            conn.close()
            self._send_json({"ok": True, "id": new_id}, status=201)
        except (json.JSONDecodeError, KeyError, ValueError) as e:
            self._send_json({"ok": False, "error": str(e)}, status=400)

    def _update_md_save(self, save_id):
        raw = self._read_body()
        try:
            payload = json.loads(raw)
            name = payload.get("name", "").strip()
            content = payload.get("content", "")
            if not name:
                self._send_json({"ok": False, "error": "이름을 입력하세요."}, status=400)
                return
            conn = get_conn()
            # Save current content as a version before updating
            old = conn.execute(
                "SELECT content FROM md_saves WHERE id=?", (save_id,)
            ).fetchone()
            if not old:
                conn.close()
                self._send_json({"ok": False, "error": "not found"}, status=404)
                return
            next_vn = conn.execute(
                "SELECT next_version FROM md_saves WHERE id=?", (save_id,)
            ).fetchone()["next_version"]
            conn.execute(
                "INSERT INTO md_versions (save_id, content, version_num) VALUES (?, ?, ?)",
                (save_id, old["content"], next_vn),
            )
            conn.execute(
                "UPDATE md_saves SET name=?, content=?, next_version=?, updated_at=CURRENT_TIMESTAMP WHERE id=?",
                (name, content, next_vn + 1, save_id),
            )
            # Keep max 30 versions only (excluding archived)
            conn.execute(
                "DELETE FROM md_versions WHERE save_id=? AND archived=0 AND id NOT IN "
                "(SELECT id FROM md_versions WHERE save_id=? AND archived=0 ORDER BY id DESC LIMIT 30)",
                (save_id, save_id),
            )
            conn.commit()
            conn.close()
            self._send_json({"ok": True})
        except (json.JSONDecodeError, KeyError, ValueError) as e:
            self._send_json({"ok": False, "error": str(e)}, status=400)

    def _delete_md_save(self, save_id):
        conn = get_conn()
        conn.execute("DELETE FROM md_versions WHERE save_id=?", (save_id,))
        cur = conn.execute("DELETE FROM md_saves WHERE id=?", (save_id,))
        conn.commit()
        deleted = cur.rowcount
        conn.close()
        if deleted == 0:
            self._send_json({"ok": False, "error": "not found"}, status=404)
            return
        self._send_json({"ok": True})

    def _list_md_versions(self, save_id):
        conn = get_conn()
        try:
            rows = conn.execute(
                "SELECT id, save_id, version_num, archived, comment, created_at FROM md_versions WHERE save_id=? ORDER BY id DESC",
                (save_id,),
            ).fetchall()
            conn.close()
            self._send_json([dict(r) for r in rows])
        except Exception as e:
            conn.close()
            self._send_json({"ok": False, "error": str(e)}, status=500)

    def _update_md_version_comment(self, version_id):
        raw = self._read_body()
        try:
            payload = json.loads(raw)
            comment = payload.get("comment", "").strip()
            conn = get_conn()
            cur = conn.execute("UPDATE md_versions SET comment=? WHERE id=?", (comment, version_id))
            conn.commit()
            conn.close()
            if cur.rowcount == 0:
                self._send_json({"ok": False, "error": "not found"}, status=404)
            else:
                self._send_json({"ok": True})
        except Exception as e:
            self._send_json({"ok": False, "error": str(e)}, status=500)

    def _toggle_md_version_archive(self, version_id):
        conn = get_conn()
        row = conn.execute("SELECT archived FROM md_versions WHERE id=?", (version_id,)).fetchone()
        if not row:
            conn.close()
            self._send_json({"ok": False, "error": "not found"}, status=404)
            return
        new_val = 0 if row["archived"] else 1
        conn.execute("UPDATE md_versions SET archived=? WHERE id=?", (new_val, version_id))
        conn.commit()
        conn.close()
        self._send_json({"ok": True, "archived": new_val})

    def _get_md_version(self, version_id):
        conn = get_conn()
        row = conn.execute(
            "SELECT * FROM md_versions WHERE id=?", (version_id,)
        ).fetchone()
        conn.close()
        if not row:
            self._send_json({"ok": False, "error": "not found"}, status=404)
            return
        self._send_json(dict(row))

    def _rollback_md_version(self, save_id, version_id):
        conn = get_conn()
        ver = conn.execute(
            "SELECT content FROM md_versions WHERE id=? AND save_id=?",
            (version_id, save_id),
        ).fetchone()
        if not ver:
            conn.close()
            self._send_json({"ok": False, "error": "version not found"}, status=404)
            return
        # Save current content as a version before rollback
        current = conn.execute(
            "SELECT content FROM md_saves WHERE id=?", (save_id,)
        ).fetchone()
        if current:
            next_vn = conn.execute(
                "SELECT next_version FROM md_saves WHERE id=?", (save_id,)
            ).fetchone()["next_version"]
            conn.execute(
                "INSERT INTO md_versions (save_id, content, version_num) VALUES (?, ?, ?)",
                (save_id, current["content"], next_vn),
            )
            conn.execute(
                "UPDATE md_saves SET next_version=? WHERE id=?",
                (next_vn + 1, save_id),
            )
            # Keep max 30 versions only (excluding archived)
            conn.execute(
                "DELETE FROM md_versions WHERE save_id=? AND archived=0 AND id NOT IN "
                "(SELECT id FROM md_versions WHERE save_id=? AND archived=0 ORDER BY id DESC LIMIT 30)",
                (save_id, save_id),
            )
        conn.execute(
            "UPDATE md_saves SET content=?, updated_at=CURRENT_TIMESTAMP WHERE id=?",
            (ver["content"], save_id),
        )
        conn.commit()
        conn.close()
        self._send_json({"ok": True, "content": ver["content"]})

    def _csv_encode(self):
        raw = self._read_body()
        try:
            payload = json.loads(raw)
            text = payload.get("text", "")
            encoding = payload.get("encoding", "utf-8").strip().lower()
            filename = payload.get("filename", "data.csv").strip()
            encoded = text.encode(encoding)
            self.send_response(200)
            self.send_header("Content-Type", "text/csv")
            self.send_header("Content-Disposition", f'attachment; filename="{filename}"')
            self.send_header("Content-Length", str(len(encoded)))
            self.end_headers()
            self.wfile.write(encoded)
        except (json.JSONDecodeError, KeyError) as e:
            self._send_json({"ok": False, "error": f"잘못된 요청: {e}"}, status=400)
        except LookupError:
            self._send_json({"ok": False, "error": f"지원하지 않는 인코딩: {encoding}"}, status=400)
        except UnicodeEncodeError as e:
            self._send_json({"ok": False, "error": f"인코딩 변환 실패: {e}"}, status=400)

    def _proxy_mock(self, method, path):
        logger.debug(f"[Mock] {method} {path}")
        body_raw = self._read_body()
        request_headers = {k: v for k, v in self.headers.items()}
        req_json = None
        if body_raw.strip():
            try:
                req_json = json.loads(body_raw)
            except json.JSONDecodeError:
                req_json = None

        conn = get_conn()
        candidates = conn.execute(
            "SELECT * FROM mocks WHERE method=? AND path=? ORDER BY id DESC",
            (method, path),
        ).fetchall()
        conn.close()

        for row in candidates:
            expected = json.loads(row["request_json"]) if row["request_json"] else None
            if expected is not None and req_json != expected:
                continue

            headers = json.loads(row["response_headers"]) if row["response_headers"] else {}
            status = row["response_status"]
            response_body = row["response_body"]
            body_obj = json.loads(response_body) if response_body else None

            body_bytes = b""
            ctype = None
            if isinstance(body_obj, (dict, list)):
                body_bytes = json.dumps(body_obj, ensure_ascii=False).encode("utf-8")
                ctype = "application/json; charset=utf-8"
            elif body_obj is None:
                body_bytes = b""
                ctype = "application/json; charset=utf-8"
            else:
                body_bytes = str(body_obj).encode("utf-8")
                ctype = "text/plain; charset=utf-8"

            self.send_response(status)
            lowered = {str(k).lower() for k in headers.keys()}
            sent_headers = dict(headers)
            if "content-type" not in lowered:
                self.send_header("Content-Type", ctype)
                sent_headers["Content-Type"] = ctype
            for k, v in headers.items():
                self.send_header(str(k), str(v))
            self.send_header("Content-Length", str(len(body_bytes)))
            sent_headers["Content-Length"] = str(len(body_bytes))
            self.end_headers()
            if body_bytes:
                self.wfile.write(body_bytes)
            self._log_traffic(
                matched=True,
                matched_mock_id=row["id"],
                method=method,
                path=path,
                request_headers=request_headers,
                request_body=body_raw,
                request_json=req_json,
                response_status=status,
                response_headers=sent_headers,
                response_body=body_obj,
            )
            return

        miss_payload = {
            "ok": False,
            "error": "No matching mock found",
            "method": method,
            "path": path,
            "received_body": req_json,
        }
        self._log_traffic(
            matched=False,
            matched_mock_id=None,
            method=method,
            path=path,
            request_headers=request_headers,
            request_body=body_raw,
            request_json=req_json,
            response_status=404,
            response_headers={"Content-Type": "application/json; charset=utf-8"},
            response_body=miss_payload,
        )
        self._send_json(
            miss_payload,
            status=404,
        )

    def do_GET(self):
        parsed = urlparse(self.path)
        path = parsed.path
        query = parse_qs(parsed.query)

        if path == "/favicon.ico":
            self.send_response(204)
            self.send_header("Content-Length", "0")
            self.end_headers()
            return

        if path == "/api/mocks":
            self._list_mocks()
            return
        if path == "/api/logs":
            limit = query.get("limit", ["200"])[0]
            self._list_logs(limit=limit)
            return
        if path == "/api/json/saves":
            self._list_json_saves()
            return
        if path.startswith("/api/json/saves/"):
            try:
                save_id = int(path.split("/")[-1])
                self._get_json_save(save_id)
                return
            except ValueError:
                self._send_json({"ok": False, "error": "invalid id"}, status=400)
                return
        if path == "/api/csv/saves":
            self._list_csv_saves()
            return
        if path.startswith("/api/csv/saves/"):
            try:
                save_id = int(path.split("/")[-1])
                self._get_csv_save(save_id)
                return
            except ValueError:
                self._send_json({"ok": False, "error": "invalid id"}, status=400)
                return
        if path == "/api/md/saves":
            self._list_md_saves()
            return
        # /api/md/saves/:id/html (HTML for popup preview)
        m_html = re.match(r"^/api/md/saves/(\d+)/html$", path)
        if m_html:
            self._get_md_save_html(int(m_html.group(1)))
            return
        # /api/md/saves/:id/versions
        m_ver = re.match(r"^/api/md/saves/(\d+)/versions$", path)
        if m_ver:
            self._list_md_versions(int(m_ver.group(1)))
            return
        # /api/md/versions/:id (single version detail)
        m_ver_detail = re.match(r"^/api/md/versions/(\d+)$", path)
        if m_ver_detail:
            self._get_md_version(int(m_ver_detail.group(1)))
            return
        if path.startswith("/api/md/saves/"):
            try:
                save_id = int(path.split("/")[-1])
                self._get_md_save(save_id)
                return
            except ValueError:
                self._send_json({"ok": False, "error": "invalid id"}, status=400)
                return
        # /api/md/proofread/result?save_id=N
        if path == "/api/md/proofread/result":
            save_id = query.get("save_id", [None])[0]
            if save_id is not None:
                self._get_proofread_result(int(save_id))
            else:
                self._send_json({"ok": False, "error": "save_id required"}, status=400)
            return

        # ── Data AI GET ──
        if path == "/api/dataai/saves":
            self._dataai_list()
            return
        m_dataai = re.match(r"^/api/dataai/saves/(\d+)$", path)
        if m_dataai:
            self._dataai_get(int(m_dataai.group(1)))
            return

        # ── Git management GET ──
        if path == "/api/git/status":
            self._git_status(query)
            return
        if path == "/api/git/branches":
            self._git_branches(query)
            return
        if path == "/api/git/diff":
            self._git_diff(query)
            return
        if path == "/api/git/log":
            self._git_log(query)
            return
        if path == "/api/git/templates":
            self._git_list_templates(query)
            return

        # ── Developer mode GET ──
        if path == "/api/dev/auth/status":
            self._dev_auth_status()
            return
        if path == "/api/dev/tables":
            self._dev_list_tables()
            return
        m_tbl = re.match(r"^/api/dev/tables/([a-zA-Z_]\w*)$", path)
        if m_tbl:
            self._dev_get_table(m_tbl.group(1))
            return
        if path == "/api/dev/tabs":
            self._dev_get_tabs()
            return
        if path == "/api/dev/modules":
            self._dev_get_modules()
            return
        if path == "/api/dev/site-config":
            self._dev_get_site_config()
            return
        if path == "/api/dev/onboarding":
            self._dev_get_onboarding()
            return
        if path == "/api/ai/providers":
            self._ai_get_providers()
            return
        if path == "/api/dev/ai-keys":
            self._dev_get_ai_keys()
            return
        if path == "/api/dev/version":
            self._dev_get_version()
            return
        if path == "/api/dev/releases":
            self._dev_get_releases()
            return
        if path == "/api/dev/cdn/status":
            self._dev_cdn_status()
            return
        if path == "/api/dev/cdn/check-latest":
            self._dev_cdn_check_latest()
            return
        # ── i18n API ──
        m_lang = re.match(r"^/api/lang/([a-z]{2})$", path)
        if m_lang:
            self._get_translations(m_lang.group(1))
            return

        # ── Custom Plugin routes (GET) ──
        if path == "/api/custom/plugins":
            self._custom_list_plugins()
            return
        m_custom = re.match(r"^/api/custom/([a-zA-Z0-9_-]+)/(.+)$", path)
        if m_custom:
            self._handle_custom_route(m_custom.group(1), m_custom.group(2), "GET")
            return

        if path.startswith("/api/"):
            self._send_json({"ok": False, "error": "unsupported api"}, status=404)
            return

        if self._serve_static(path):
            return
        self._proxy_mock("GET", path)

    def do_POST(self):
        parsed = urlparse(self.path)
        path = parsed.path

        if path == "/api/mocks":
            self._create_mock()
            return
        if path == "/api/translate":
            self._translate()
            return
        if path == "/api/csv/encode":
            self._csv_encode()
            return
        if path == "/api/json/saves":
            self._create_json_save()
            return
        if path == "/api/csv/saves":
            self._create_csv_save()
            return
        if path == "/api/md/saves":
            self._create_md_save()
            return
        if path == "/api/md/proofread":
            self._md_proofread()
            return
        if path == "/api/md/proofread/save":
            self._save_proofread_result()
            return
        if path == "/api/dataai/generate":
            self._dataai_generate()
            return
        if path == "/api/dataai/saves":
            self._dataai_save()
            return
        # /api/md/saves/:id/rollback/:versionId
        m_rb = re.match(r"^/api/md/saves/(\d+)/rollback/(\d+)$", path)
        if m_rb:
            self._rollback_md_version(int(m_rb.group(1)), int(m_rb.group(2)))
            return
        # ── Developer mode POST ──
        if path == "/api/dev/auth/register":
            self._dev_auth_register()
            return
        if path == "/api/dev/auth/login":
            self._dev_auth_login()
            return
        if path == "/api/dev/query":
            self._dev_query()
            return
        if path == "/api/dev/cdn/sync":
            self._dev_cdn_sync()
            return
        if path == "/api/dev/update":
            self._dev_update()
            return
        # ── PDF conversion ──
        if path == "/api/pdf/convert/xlsx":
            self._pdf_convert_xlsx()
            return
        if path == "/api/pdf/convert/pptx":
            self._pdf_convert_pptx()
            return

        # ── Git management POST ──
        if path == "/api/git/pick-repo":
            self._diff_pick_folder()
            return
        if path == "/api/git/commit":
            self._git_commit()
            return
        if path == "/api/git/discard":
            self._git_discard()
            return
        if path == "/api/git/switch-branch":
            self._git_switch_branch()
            return
        if path == "/api/git/create-branch":
            self._git_create_branch()
            return
        if path == "/api/git/templates":
            self._git_create_template()
            return

        if path == "/api/diff/folder":
            self._diff_compare_folder()
            return

        if path == "/api/diff/file":
            self._diff_read_file()
            return

        if path == "/api/diff/pick-file":
            self._diff_pick_file()
            return

        if path == "/api/diff/pick-folder":
            self._diff_pick_folder()
            return

        if path == "/api/diff/save":
            self._diff_save_file()
            return

        if path == "/api/dev/onboarding/complete":
            self._dev_complete_onboarding()
            return

        # ── Custom Plugin routes (POST) ──
        m_custom = re.match(r"^/api/custom/([a-zA-Z0-9_-]+)/(.+)$", path)
        if m_custom:
            self._handle_custom_route(m_custom.group(1), m_custom.group(2), "POST")
            return

        self._proxy_mock("POST", path)

    def do_PUT(self):
        parsed = urlparse(self.path)
        path = parsed.path
        if path.startswith("/api/mocks/"):
            try:
                mock_id = int(path.split("/")[-1])
                self._update_mock(mock_id)
                return
            except ValueError:
                self._send_json({"ok": False, "error": "invalid id"}, status=400)
                return
        if path.startswith("/api/json/saves/"):
            try:
                save_id = int(path.split("/")[-1])
                self._update_json_save(save_id)
                return
            except ValueError:
                self._send_json({"ok": False, "error": "invalid id"}, status=400)
                return
        if path.startswith("/api/csv/saves/"):
            try:
                save_id = int(path.split("/")[-1])
                self._update_csv_save(save_id)
                return
            except ValueError:
                self._send_json({"ok": False, "error": "invalid id"}, status=400)
                return
        # /api/md/versions/:id/comment (edit comment)
        m_comment = re.match(r"^/api/md/versions/(\d+)/comment$", path)
        if m_comment:
            self._update_md_version_comment(int(m_comment.group(1)))
            return
        # /api/md/versions/:id/archive (toggle archive)
        m_arch = re.match(r"^/api/md/versions/(\d+)/archive$", path)
        if m_arch:
            self._toggle_md_version_archive(int(m_arch.group(1)))
            return
        if path.startswith("/api/md/saves/"):
            try:
                save_id = int(path.split("/")[-1])
                self._update_md_save(save_id)
                return
            except ValueError:
                self._send_json({"ok": False, "error": "invalid id"}, status=400)
                return
        # ── Developer mode PUT ──
        m_dev_row = re.match(r"^/api/dev/tables/([a-zA-Z_]\w*)/(.+)$", path)
        if m_dev_row:
            row_id = unquote(m_dev_row.group(2))
            try:
                row_id = int(row_id)
            except ValueError:
                pass
            self._dev_update_row(m_dev_row.group(1), row_id)
            return
        if path == "/api/dev/tabs":
            self._dev_save_tabs()
            return
        if path == "/api/dev/modules":
            self._dev_save_modules()
            return
        if path == "/api/dev/site-config":
            self._dev_save_site_config()
            return
        if path == "/api/dev/ai-keys":
            self._dev_save_ai_keys()
            return
        # ── Custom Plugin routes (PUT) ──
        if path == "/api/custom/plugins/toggle":
            self._custom_toggle_plugin()
            return
        m_custom = re.match(r"^/api/custom/([a-zA-Z0-9_-]+)/(.+)$", path)
        if m_custom:
            self._handle_custom_route(m_custom.group(1), m_custom.group(2), "PUT")
            return
        self._proxy_mock("PUT", path)

    def do_DELETE(self):
        parsed = urlparse(self.path)
        path = parsed.path
        if path.startswith("/api/mocks/"):
            try:
                mock_id = int(path.split("/")[-1])
                self._delete_mock(mock_id)
                return
            except ValueError:
                self._send_json({"ok": False, "error": "invalid id"}, status=400)
                return
        if path == "/api/logs":
            self._clear_logs()
            return
        if path.startswith("/api/json/saves/"):
            try:
                save_id = int(path.split("/")[-1])
                self._delete_json_save(save_id)
                return
            except ValueError:
                self._send_json({"ok": False, "error": "invalid id"}, status=400)
                return
        if path.startswith("/api/csv/saves/"):
            try:
                save_id = int(path.split("/")[-1])
                self._delete_csv_save(save_id)
                return
            except ValueError:
                self._send_json({"ok": False, "error": "invalid id"}, status=400)
                return
        if path.startswith("/api/md/saves/"):
            try:
                save_id = int(path.split("/")[-1])
                self._delete_md_save(save_id)
                return
            except ValueError:
                self._send_json({"ok": False, "error": "invalid id"}, status=400)
                return
        m_pr = re.match(r"^/api/md/proofread/result/(\d+)$", path)
        if m_pr:
            self._delete_proofread_result(int(m_pr.group(1)))
            return
        m_dataai_del = re.match(r"^/api/dataai/saves/(\d+)$", path)
        if m_dataai_del:
            self._dataai_delete(int(m_dataai_del.group(1)))
            return
        # ── Git template DELETE ──
        m_git_tpl = re.match(r"^/api/git/templates/(\d+)$", path)
        if m_git_tpl:
            self._git_delete_template(int(m_git_tpl.group(1)))
            return
        # ── Developer mode DELETE ──
        m_dev_del = re.match(r"^/api/dev/tables/([a-zA-Z_]\w*)/(.+)$", path)
        if m_dev_del:
            row_id = unquote(m_dev_del.group(2))
            try:
                row_id = int(row_id)
            except ValueError:
                pass
            self._dev_delete_row(m_dev_del.group(1), row_id)
            return
        # ── Custom Plugin routes (DELETE) ──
        m_custom = re.match(r"^/api/custom/([a-zA-Z0-9_-]+)/(.+)$", path)
        if m_custom:
            self._handle_custom_route(m_custom.group(1), m_custom.group(2), "DELETE")
            return
        self._proxy_mock("DELETE", path)

    def do_PATCH(self):
        parsed = urlparse(self.path)
        path = parsed.path
        self._proxy_mock("PATCH", path)

    def do_OPTIONS(self):
        parsed = urlparse(self.path)
        path = parsed.path
        self._proxy_mock("OPTIONS", path)

    def do_HEAD(self):
        parsed = urlparse(self.path)
        path = parsed.path
        self._proxy_mock("HEAD", path)

    # ── Developer mode API ──

    def _dev_check_auth(self):
        """Verify auth token. Returns False after sending 401 on failure."""
        token = self.headers.get("X-Dev-Token", "")
        if token and token in _dev_sessions:
            return True
        self._send_json({"ok": False, "error": "authentication required"}, status=401)
        return False

    def _dev_auth_status(self):
        conn = get_conn()
        row = conn.execute("SELECT value FROM dev_settings WHERE key='auth_username'").fetchone()
        conn.close()
        self._send_json({"ok": True, "registered": row is not None})

    def _dev_auth_register(self):
        body = self._read_body()
        data = json.loads(body)
        username = data.get("username", "").strip()
        password = data.get("password", "").strip()
        if not username or not password:
            self._send_json({"ok": False, "error": "username and password required"}, status=400)
            return
        conn = get_conn()
        existing = conn.execute("SELECT value FROM dev_settings WHERE key='auth_username'").fetchone()
        if existing:
            conn.close()
            self._send_json({"ok": False, "error": "account already exists"}, status=409)
            return
        salt = secrets.token_bytes(32)
        pw_hash = _hash_password(password, salt)
        conn.execute("INSERT OR REPLACE INTO dev_settings (key, value) VALUES ('auth_username', ?)", (username,))
        conn.execute("INSERT OR REPLACE INTO dev_settings (key, value) VALUES ('auth_password_hash', ?)", (pw_hash,))
        conn.execute("INSERT OR REPLACE INTO dev_settings (key, value) VALUES ('auth_salt', ?)", (salt.hex(),))
        conn.commit()
        conn.close()
        token = secrets.token_hex(32)
        _dev_sessions[token] = True
        self._send_json({"ok": True, "token": token})

    def _dev_auth_login(self):
        body = self._read_body()
        data = json.loads(body)
        username = data.get("username", "").strip()
        password = data.get("password", "").strip()
        conn = get_conn()
        stored_user = conn.execute("SELECT value FROM dev_settings WHERE key='auth_username'").fetchone()
        stored_hash = conn.execute("SELECT value FROM dev_settings WHERE key='auth_password_hash'").fetchone()
        stored_salt = conn.execute("SELECT value FROM dev_settings WHERE key='auth_salt'").fetchone()
        conn.close()
        if not stored_user or not stored_hash or not stored_salt:
            self._send_json({"ok": False, "error": "no account registered"}, status=404)
            return
        if username != stored_user[0]:
            self._send_json({"ok": False, "error": "invalid credentials"}, status=401)
            return
        salt = bytes.fromhex(stored_salt[0])
        if _hash_password(password, salt) != stored_hash[0]:
            self._send_json({"ok": False, "error": "invalid credentials"}, status=401)
            return
        token = secrets.token_hex(32)
        _dev_sessions[token] = True
        self._send_json({"ok": True, "token": token})

    def _dev_list_tables(self):
        if not self._dev_check_auth():
            return
        conn = get_conn()
        rows = conn.execute("SELECT name FROM sqlite_master WHERE type='table' AND name NOT LIKE 'sqlite_%' ORDER BY name").fetchall()
        conn.close()
        self._send_json({"ok": True, "tables": [r[0] for r in rows]})

    def _dev_get_table(self, table_name):
        if not self._dev_check_auth():
            return
        conn = get_conn()
        # Verify table exists
        exists = conn.execute("SELECT name FROM sqlite_master WHERE type='table' AND name=?", (table_name,)).fetchone()
        if not exists:
            conn.close()
            self._send_json({"ok": False, "error": "table not found"}, status=404)
            return
        # Column info
        columns = conn.execute(f"PRAGMA table_info([{table_name}])").fetchall()
        col_info = [{"cid": c[0], "name": c[1], "type": c[2], "notnull": c[3], "default": c[4], "pk": c[5]} for c in columns]
        # Data
        rows = conn.execute(f"SELECT * FROM [{table_name}] ORDER BY rowid DESC LIMIT 500").fetchall()
        col_names = [c[1] for c in columns]
        data = [dict(zip(col_names, row)) for row in rows]
        conn.close()
        self._send_json({"ok": True, "columns": col_info, "data": data})

    def _dev_update_row(self, table_name, row_id):
        if not self._dev_check_auth():
            return
        body = self._read_body()
        data = json.loads(body)
        conn = get_conn()
        try:
            exists = conn.execute("SELECT name FROM sqlite_master WHERE type='table' AND name=?", (table_name,)).fetchone()
            if not exists:
                conn.close()
                self._send_json({"ok": False, "error": "table not found"}, status=404)
                return
            # Find PK column
            columns = conn.execute(f"PRAGMA table_info([{table_name}])").fetchall()
            pk_col = next((c[1] for c in columns if c[5] == 1), "id")
            sets = ", ".join(f"[{k}] = ?" for k in data.keys())
            # Convert dict/list values to JSON strings for SQLite binding
            vals = [json.dumps(v, ensure_ascii=False) if isinstance(v, (dict, list)) else v for v in data.values()]
            vals.append(row_id)
            conn.execute(f"UPDATE [{table_name}] SET {sets} WHERE [{pk_col}] = ?", vals)
            conn.commit()
            self._send_json({"ok": True})
        except Exception as e:
            logger.error("[dev] update row error: %s", e)
            self._send_json({"ok": False, "error": str(e)}, status=500)
        finally:
            conn.close()

    def _dev_delete_row(self, table_name, row_id):
        if not self._dev_check_auth():
            return
        conn = get_conn()
        exists = conn.execute("SELECT name FROM sqlite_master WHERE type='table' AND name=?", (table_name,)).fetchone()
        if not exists:
            conn.close()
            self._send_json({"ok": False, "error": "table not found"}, status=404)
            return
        columns = conn.execute(f"PRAGMA table_info([{table_name}])").fetchall()
        pk_col = next((c[1] for c in columns if c[5] == 1), "id")
        conn.execute(f"DELETE FROM [{table_name}] WHERE [{pk_col}] = ?", (row_id,))
        conn.commit()
        conn.close()
        self._send_json({"ok": True})

    def _dev_query(self):
        if not self._dev_check_auth():
            return
        body = self._read_body()
        data = json.loads(body)
        sql = data.get("sql", "").strip()
        if not sql:
            self._send_json({"ok": False, "error": "sql required"}, status=400)
            return
        # Allow only SELECT/PRAGMA, block multiple statements
        first_word = sql.split()[0].upper() if sql.split() else ""
        if first_word not in ("SELECT", "PRAGMA"):
            self._send_json({"ok": False, "error": "only SELECT/PRAGMA allowed"}, status=403)
            return
        # Strip trailing semicolons/whitespace then check for remaining semicolons (multi-statement)
        stripped = sql.strip().rstrip(";").strip()
        if ";" in stripped:
            self._send_json({"ok": False, "error": "multiple statements not allowed"}, status=403)
            return
        conn = get_conn()
        try:
            # Unset row_factory to get tuples (preserves duplicate column names)
            conn.row_factory = None
            cursor = conn.execute(stripped)
            raw_names = [d[0] for d in cursor.description] if cursor.description else []
            # Append suffix to duplicate column names (id, id → id, id_2)
            col_names = []
            seen = {}
            for name in raw_names:
                if name in seen:
                    seen[name] += 1
                    col_names.append(f"{name}_{seen[name]}")
                else:
                    seen[name] = 1
                    col_names.append(name)
            rows = cursor.fetchall()
            data_out = [dict(zip(col_names, row)) for row in rows]
            conn.row_factory = sqlite3.Row
            conn.close()
            self._send_json({"ok": True, "columns": col_names, "data": data_out})
        except Exception as e:
            conn.close()
            self._send_json({"ok": False, "error": str(e)}, status=400)

    def _dev_get_tabs(self):
        conn = get_conn()
        row = conn.execute("SELECT value FROM dev_settings WHERE key='tabs_config'").fetchone()
        conn.close()
        default_tabs = [
            {"id": "mock", "label": "Mock Server", "visible": True, "order": 0},
            {"id": "charcount", "label": "Char Count", "visible": True, "order": 1},
            {"id": "mybatis", "label": "MyBatis Log", "visible": True, "order": 2},
            {"id": "jsonformat", "label": "JSON Format", "visible": True, "order": 3},
            {"id": "translate", "label": "Translate", "visible": True, "order": 4},
            {"id": "csv", "label": "CSV Editor", "visible": True, "order": 5},
            {"id": "markdown", "label": "Markdown", "visible": True, "order": 6},
            {"id": "paramchanger", "label": "Param Convert", "visible": True, "order": 7},
            {"id": "pdfconvert", "label": "PDF Convert", "visible": True, "order": 8},
            {"id": "diffcompare", "label": "Diff Compare", "visible": True, "order": 9},
            {"id": "git", "label": "Git", "visible": True, "order": 10},
        ]
        if row:
            saved = json.loads(row[0])
            saved_ids = {t["id"] for t in saved}
            max_order = max((t.get("order", 0) for t in saved), default=0)
            for dt in default_tabs:
                if dt["id"] not in saved_ids:
                    max_order += 1
                    dt["order"] = max_order
                    saved.append(dt)
            self._send_json({"ok": True, "tabs": saved})
        else:
            self._send_json({"ok": True, "tabs": default_tabs})

    def _dev_save_tabs(self):
        body = self._read_body()
        data = json.loads(body)
        tabs = data.get("tabs", [])
        conn = get_conn()
        conn.execute(
            "INSERT OR REPLACE INTO dev_settings (key, value, encrypted) VALUES ('tabs_config', ?, 0)",
            (json.dumps(tabs, ensure_ascii=False),),
        )
        conn.commit()
        conn.close()
        self._send_json({"ok": True})

    # ══════════════════════════════════
    # Diff comparison API
    # ══════════════════════════════════
    def _diff_pick_file(self):
        """Native file picker dialog via osascript"""
        import subprocess
        try:
            result = subprocess.run(
                ["osascript", "-e", 'POSIX path of (choose file with prompt "파일 선택")'],
                capture_output=True, text=True, timeout=120
            )
            path = result.stdout.strip()
            if path:
                self._send_json({"ok": True, "path": path})
            else:
                self._send_json({"ok": False, "error": "파일을 선택하지 않았습니다"})
        except subprocess.TimeoutExpired:
            self._send_json({"ok": False, "error": "시간 초과"})
        except Exception as e:
            self._send_json({"ok": False, "error": str(e)})

    def _diff_pick_folder(self):
        """Native folder picker dialog via osascript"""
        import subprocess
        try:
            result = subprocess.run(
                ["osascript", "-e", 'POSIX path of (choose folder with prompt "폴더 선택")'],
                capture_output=True, text=True, timeout=120
            )
            path = result.stdout.strip().rstrip("/")
            if path:
                self._send_json({"ok": True, "path": path})
            else:
                self._send_json({"ok": False, "error": "폴더를 선택하지 않았습니다"})
        except subprocess.TimeoutExpired:
            self._send_json({"ok": False, "error": "시간 초과"})
        except Exception as e:
            self._send_json({"ok": False, "error": str(e)})

    def _diff_save_file(self):
        """Save content to file path"""
        import os
        body = json.loads(self._read_body())
        file_path = body.get("path", "").strip()
        content = body.get("content", "")
        if not file_path:
            self._send_json({"ok": False, "error": "경로가 비어있습니다"})
            return
        file_path = os.path.expanduser(file_path)
        if not os.path.isfile(file_path):
            self._send_json({"ok": False, "error": "파일을 찾을 수 없습니다: " + file_path})
            return
        try:
            with open(file_path, "w", encoding="utf-8") as f:
                f.write(content)
            self._send_json({"ok": True})
        except Exception as e:
            self._send_json({"ok": False, "error": str(e)})

    def _diff_read_file(self):
        """Read file content from given path"""
        import os
        body = json.loads(self._read_body())
        file_path = body.get("path", "").strip()
        if not file_path:
            self._send_json({"ok": False, "error": "경로가 비어있습니다"})
            return
        file_path = os.path.expanduser(file_path)
        if not os.path.isfile(file_path):
            self._send_json({"ok": False, "error": "파일을 찾을 수 없습니다: " + file_path})
            return
        size = os.path.getsize(file_path)
        if size > 10 * 1024 * 1024:
            self._send_json({"ok": False, "error": "파일이 10MB를 초과합니다"})
            return
        try:
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                content = f.read()
            self._send_json({"ok": True, "content": content, "name": os.path.basename(file_path)})
        except Exception as e:
            self._send_json({"ok": False, "error": str(e)})

    def _diff_compare_folder(self):
        """Recursively compare files in two folders and return results"""
        import os

        body = json.loads(self._read_body())
        left_path = body.get("left", "")
        right_path = body.get("right", "")

        if not os.path.isdir(left_path):
            self._send_json({"ok": False, "error": f"왼쪽 폴더가 존재하지 않습니다: {left_path}"}, 400)
            return
        if not os.path.isdir(right_path):
            self._send_json({"ok": False, "error": f"오른쪽 폴더가 존재하지 않습니다: {right_path}"}, 400)
            return

        MAX_SIZE = 10 * 1024 * 1024  # 10MB
        TEXT_EXTENSIONS = {
            ".txt", ".md", ".py", ".js", ".ts", ".jsx", ".tsx", ".html", ".htm",
            ".css", ".scss", ".less", ".json", ".xml", ".yaml", ".yml", ".toml",
            ".ini", ".cfg", ".conf", ".sh", ".bash", ".zsh", ".bat", ".cmd",
            ".java", ".c", ".cpp", ".h", ".hpp", ".cs", ".go", ".rs", ".rb",
            ".php", ".pl", ".sql", ".r", ".m", ".swift", ".kt", ".scala",
            ".vue", ".svelte", ".astro", ".env", ".gitignore", ".dockerfile",
            ".csv", ".tsv", ".log", ".properties",
        }

        def is_text_file(path):
            _, ext = os.path.splitext(path)
            return ext.lower() in TEXT_EXTENSIONS

        def get_files(base):
            result = {}
            for root, dirs, files in os.walk(base):
                dirs[:] = [d for d in dirs if not d.startswith(".")]
                for f in files:
                    if f.startswith("."):
                        continue
                    full = os.path.join(root, f)
                    rel = os.path.relpath(full, base)
                    if is_text_file(f) and os.path.getsize(full) <= MAX_SIZE:
                        result[rel] = full
            return result

        left_files = get_files(left_path)
        right_files = get_files(right_path)
        all_paths = sorted(set(left_files.keys()) | set(right_files.keys()))

        files = []
        for rel in all_paths:
            left_full = left_files.get(rel)
            right_full = right_files.get(rel)

            left_content = ""
            right_content = ""
            status = "same"

            if left_full and right_full:
                try:
                    left_content = open(left_full, "r", encoding="utf-8", errors="replace").read()
                    right_content = open(right_full, "r", encoding="utf-8", errors="replace").read()
                except Exception:
                    continue
                status = "same" if left_content == right_content else "modified"
            elif left_full:
                try:
                    left_content = open(left_full, "r", encoding="utf-8", errors="replace").read()
                except Exception:
                    continue
                status = "removed"
            else:
                try:
                    right_content = open(right_full, "r", encoding="utf-8", errors="replace").read()
                except Exception:
                    continue
                status = "added"

            files.append({
                "path": rel,
                "status": status,
                "left": left_content,
                "right": right_content,
            })

        self._send_json({"ok": True, "files": files})

    # ══════════════════════════════════
    # Git management API
    # ══════════════════════════════════
    def _git_run(self, repo_path, args):
        """Execute git command, returns (ok, stdout, stderr)"""
        result = subprocess.run(
            ["git"] + args, cwd=repo_path,
            capture_output=True, text=True, timeout=30
        )
        return result.returncode == 0, result.stdout, result.stderr.strip()

    def _git_validate_repo(self, path):
        """Validate git repository. Returns True if valid, sends error response and returns False otherwise"""
        if not path or not os.path.isdir(path):
            self._send_json({"ok": False, "error": "유효하지 않은 경로입니다"}, 400)
            return False
        ok, _, _ = self._git_run(path, ["rev-parse", "--git-dir"])
        if not ok:
            self._send_json({"ok": False, "error": "Git 저장소가 아닙니다: " + path}, 400)
            return False
        return True

    def _git_status(self, query):
        repo = query.get("repo", [""])[0]
        if not self._git_validate_repo(repo):
            return
        # Current branch
        ok, branch, _ = self._git_run(repo, ["branch", "--show-current"])
        branch = branch.strip() if ok else "(detached)"
        # status --porcelain
        ok, out, err = self._git_run(repo, ["status", "--porcelain=v1"])
        if not ok:
            self._send_json({"ok": False, "error": err}, 500)
            return
        files = []
        for line in out.splitlines():
            if len(line) < 4:
                continue
            index_st = line[0]
            work_st = line[1]
            fname = line[3:]
            # Skip directories (untracked dirs end with "/")
            if fname.endswith("/"):
                continue
            # renamed: "R  old -> new"
            old_name = None
            if index_st == "R" or work_st == "R":
                parts = fname.split(" -> ")
                if len(parts) == 2:
                    old_name = parts[0]
                    fname = parts[1]
            # Determine status
            if index_st == "?" and work_st == "?":
                status = "untracked"
            elif index_st == "A":
                status = "added"
            elif index_st == "D" or work_st == "D":
                status = "deleted"
            elif index_st == "R":
                status = "renamed"
            elif index_st != " " and index_st != "?":
                status = "modified"
            elif work_st == "M":
                status = "modified"
            else:
                status = "modified"
            entry = {"file": fname, "status": status, "index": index_st, "worktree": work_st}
            if old_name:
                entry["old_name"] = old_name
            files.append(entry)
        # git config user.name / user.email
        _, user_name, _ = self._git_run(repo, ["config", "user.name"])
        _, user_email, _ = self._git_run(repo, ["config", "user.email"])
        self._send_json({"ok": True, "branch": branch, "files": files,
                         "user": user_name.strip(), "email": user_email.strip()})

    def _git_branches(self, query):
        repo = query.get("repo", [""])[0]
        if not self._git_validate_repo(repo):
            return
        ok, out, err = self._git_run(repo, ["branch", "-a"])
        if not ok:
            self._send_json({"ok": False, "error": err}, 500)
            return
        current = ""
        local = []
        remote = []
        for line in out.splitlines():
            line = line.strip()
            if not line:
                continue
            is_current = line.startswith("* ")
            name = line[2:] if is_current else line
            if is_current:
                current = name
            if name.startswith("remotes/"):
                remote.append(name.replace("remotes/", "", 1))
            else:
                local.append(name)
        self._send_json({"ok": True, "current": current, "local": local, "remote": remote})

    def _git_diff(self, query):
        repo = query.get("repo", [""])[0]
        file_path = query.get("file", [""])[0]
        if not self._git_validate_repo(repo):
            return
        if not file_path:
            self._send_json({"ok": False, "error": "file 파라미터가 필요합니다"}, 400)
            return
        # Binary check
        ok, numstat, _ = self._git_run(repo, ["diff", "--numstat", "--", file_path])
        if ok and numstat.strip().startswith("-\t-"):
            self._send_json({"ok": True, "binary": True, "left": "", "right": ""})
            return
        # left = HEAD version
        left = ""
        ok_head, head_content, _ = self._git_run(repo, ["show", "HEAD:" + file_path])
        if ok_head:
            left = head_content
        # right = working directory version
        right = ""
        full_path = os.path.join(repo, file_path)
        if os.path.isfile(full_path):
            try:
                with open(full_path, "r", encoding="utf-8", errors="replace") as f:
                    content = f.read()
                if len(content) > 1024 * 1024:
                    self._send_json({"ok": True, "binary": False, "left": left, "right": "(파일이 1MB를 초과합니다)", "too_large": True})
                    return
                right = content
            except Exception:
                right = ""
        self._send_json({"ok": True, "binary": False, "left": left, "right": right})

    def _git_log(self, query):
        repo = query.get("repo", [""])[0]
        limit = query.get("limit", ["20"])[0]
        if not self._git_validate_repo(repo):
            return
        try:
            limit = int(limit)
        except ValueError:
            limit = 20
        ok, out, err = self._git_run(repo, [
            "log", f"-n{limit}",
            "--pretty=format:%h\t%s\t%an\t%ai"
        ])
        if not ok:
            # Could be an empty repository
            self._send_json({"ok": True, "logs": []})
            return
        logs = []
        for line in out.splitlines():
            parts = line.split("\t", 3)
            if len(parts) == 4:
                logs.append({"hash": parts[0], "message": parts[1], "author": parts[2], "date": parts[3][:10]})
        self._send_json({"ok": True, "logs": logs})

    def _git_commit(self):
        body = json.loads(self._read_body())
        repo = body.get("repo", "")
        files = body.get("files", [])
        message = body.get("message", "").strip()
        if not self._git_validate_repo(repo):
            return
        if not files:
            self._send_json({"ok": False, "error": "커밋할 파일을 선택하세요"}, 400)
            return
        if not message:
            self._send_json({"ok": False, "error": "커밋 메시지를 입력하세요"}, 400)
            return
        logger.info(f"[Git] Commit: repo={repo}, files={len(files)}, msg={message!r}")
        # reset staging area → add selected files → commit
        self._git_run(repo, ["reset", "HEAD"])
        ok, _, err = self._git_run(repo, ["add", "--"] + files)
        if not ok:
            logger.error(f"[Git] Stage failed: {err}")
            self._send_json({"ok": False, "error": "Stage failed: " + err}, 500)
            return
        ok, out, err = self._git_run(repo, ["commit", "-m", message])
        if not ok:
            logger.error(f"[Git] Commit failed: {err}")
            self._send_json({"ok": False, "error": "Commit failed: " + err}, 500)
            return
        logger.info("[Git] Commit complete")
        self._send_json({"ok": True, "message": out.strip()})

    def _git_discard(self):
        body = json.loads(self._read_body())
        repo = body.get("repo", "")
        files = body.get("files", [])
        if not self._git_validate_repo(repo):
            return
        if not files:
            self._send_json({"ok": False, "error": "파일을 선택하세요"}, 400)
            return
        errors = []
        for f in files:
            info = f if isinstance(f, dict) else {"file": f, "status": "modified"}
            fname = info.get("file", "")
            status = info.get("status", "modified")
            if status == "untracked":
                ok, _, err = self._git_run(repo, ["clean", "-f", "--", fname])
            else:
                ok, _, err = self._git_run(repo, ["checkout", "--", fname])
            if not ok:
                errors.append(f"{fname}: {err}")
        if errors:
            self._send_json({"ok": False, "error": "\n".join(errors)}, 500)
        else:
            self._send_json({"ok": True})

    def _git_switch_branch(self):
        body = json.loads(self._read_body())
        repo = body.get("repo", "")
        branch = body.get("branch", "").strip()
        if not self._git_validate_repo(repo):
            return
        if not branch:
            self._send_json({"ok": False, "error": "브랜치명을 입력하세요"}, 400)
            return
        ok, out, err = self._git_run(repo, ["checkout", branch])
        if not ok:
            self._send_json({"ok": False, "error": err}, 500)
            return
        self._send_json({"ok": True})

    def _git_create_branch(self):
        body = json.loads(self._read_body())
        repo = body.get("repo", "")
        branch = body.get("branch", "").strip()
        if not self._git_validate_repo(repo):
            return
        if not branch:
            self._send_json({"ok": False, "error": "브랜치명을 입력하세요"}, 400)
            return
        ok, out, err = self._git_run(repo, ["checkout", "-b", branch])
        if not ok:
            self._send_json({"ok": False, "error": err}, 500)
            return
        self._send_json({"ok": True})

    # ── Git template CRUD ──
    def _git_list_templates(self, query):
        repo = query.get("repo", [""])[0]
        conn = get_conn()
        rows = conn.execute(
            "SELECT id, repo_path, name, template, created_at FROM git_commit_templates WHERE repo_path=? ORDER BY name",
            (repo,)
        ).fetchall()
        conn.close()
        self._send_json({"ok": True, "templates": [dict(r) for r in rows]})

    def _git_create_template(self):
        body = json.loads(self._read_body())
        repo = body.get("repo", "")
        name = body.get("name", "").strip()
        template = body.get("template", "")
        if not name:
            self._send_json({"ok": False, "error": "템플릿 이름을 입력하세요"}, 400)
            return
        conn = get_conn()
        conn.execute(
            "INSERT INTO git_commit_templates (repo_path, name, template) VALUES (?, ?, ?)",
            (repo, name, template)
        )
        conn.commit()
        conn.close()
        self._send_json({"ok": True})

    def _git_delete_template(self, template_id):
        conn = get_conn()
        conn.execute("DELETE FROM git_commit_templates WHERE id=?", (template_id,))
        conn.commit()
        conn.close()
        self._send_json({"ok": True})

    # ══════════════════════════════════
    # PDF conversion API
    # ══════════════════════════════════
    def _pdf_convert_xlsx(self):
        try:
            _auto_install("openpyxl")
            import openpyxl
            file_data = self._parse_multipart()
            if not file_data:
                self._send_json({"ok": False, "error": "파일 데이터 없음"}, 400)
                return
            wb = openpyxl.load_workbook(io.BytesIO(file_data), data_only=True)
            html_parts = []
            for sheet in wb.worksheets:
                html = "<h3>" + sheet.title + "</h3>"
                html += "<table style='border-collapse:collapse;width:100%;font-size:13px;margin-bottom:16px'>"
                for ri, row in enumerate(sheet.iter_rows(values_only=True)):
                    tag = "th" if ri == 0 else "td"
                    style = "border:1px solid #ddd;padding:6px 10px;text-align:left"
                    if ri == 0:
                        style += ";background:#f3f5f9;font-weight:600"
                    cells = "".join(f"<{tag} style='{style}'>{cell if cell is not None else ''}</{tag}>" for cell in row)
                    html += f"<tr>{cells}</tr>"
                html += "</table>"
                html_parts.append(html)
            self._send_json({"ok": True, "html": "\n".join(html_parts)})
        except Exception as e:
            self._send_json({"ok": False, "error": str(e)}, 500)

    def _pdf_convert_pptx(self):
        try:
            _auto_install("python-pptx")
            from pptx import Presentation
            file_data = self._parse_multipart()
            if not file_data:
                self._send_json({"ok": False, "error": "파일 데이터 없음"}, 400)
                return
            prs = Presentation(io.BytesIO(file_data))
            html_parts = []
            for i, slide in enumerate(prs.slides):
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                texts.append(text)
                html = f"<div style='border:1px solid #ddd;border-radius:8px;padding:24px;margin-bottom:16px;background:#fafbfc'>"
                html += f"<h4 style='margin:0 0 12px;color:#666'>슬라이드 {i+1}</h4>"
                html += "<br>".join(texts) if texts else "<p style='color:#999'>내용 없음</p>"
                html += "</div>"
                html_parts.append(html)
            self._send_json({"ok": True, "html": "\n".join(html_parts)})
        except Exception as e:
            self._send_json({"ok": False, "error": str(e)}, 500)

    def _get_translations(self, lang_code):
        """Return all translations for a given language code."""
        conn = get_conn()
        rows = conn.execute("SELECT key, value FROM i18n WHERE lang=?", (lang_code,)).fetchall()
        conn.close()
        translations = {r["key"]: r["value"] for r in rows}
        self._send_json({"ok": True, "lang": lang_code, "translations": translations})

    def _dev_get_site_config(self):
        conn = get_conn()
        row = conn.execute("SELECT value FROM dev_settings WHERE key='site_config'").fetchone()
        conn.close()
        if row:
            self._send_json({"ok": True, "config": json.loads(row[0])})
        else:
            self._send_json({"ok": True, "config": {"siteName": ""}})

    def _dev_save_site_config(self):
        body = self._read_body()
        data = json.loads(body)
        site_name = data.get("siteName", "").strip()
        lang = data.get("lang", "").strip() or "ko"
        toast_config = data.get("toast_config", None)
        config = {"siteName": site_name, "lang": lang}
        if toast_config:
            config["toast_config"] = toast_config
        conn = get_conn()
        conn.execute(
            "INSERT OR REPLACE INTO dev_settings (key, value, encrypted) VALUES ('site_config', ?, 0)",
            (json.dumps(config, ensure_ascii=False),),
        )
        conn.commit()
        conn.close()
        logger.info("Site config changed: name=%s, lang=%s", site_name, lang)
        self._send_json({"ok": True})

    def _dev_get_onboarding(self):
        conn = get_conn()
        row = conn.execute("SELECT value FROM dev_settings WHERE key='onboarding_complete'").fetchone()
        conn.close()
        completed = False
        if row:
            try:
                completed = json.loads(row[0]).get("completed", False)
            except Exception:
                completed = False
        self._send_json({"ok": True, "completed": completed})

    def _dev_complete_onboarding(self):
        import datetime
        config = {"completed": True, "completedAt": datetime.datetime.now().isoformat()}
        conn = get_conn()
        conn.execute(
            "INSERT OR REPLACE INTO dev_settings (key, value, encrypted) VALUES ('onboarding_complete', ?, 0)",
            (json.dumps(config),),
        )
        conn.commit()
        conn.close()
        self._send_json({"ok": True})

    # ── Custom Plugin Handlers ──

    def _custom_list_plugins(self):
        """Return list of discovered plugins with enabled state."""
        enabled_state = _get_plugin_enabled_state()
        plugins = []
        for pid, p in _custom_plugins.items():
            m = p["manifest"]
            # Default to enabled if not explicitly set
            is_enabled = enabled_state.get(pid, True)
            plugins.append({
                "id": pid,
                "name": m.get("name", pid),
                "version": m.get("version", "0.0.0"),
                "icon": m.get("icon", "puzzle"),
                "description": m.get("description", ""),
                "author": m.get("author", ""),
                "enabled": is_enabled,
                "has_routes": m.get("has_routes", False),
            })
        self._send_json({"ok": True, "plugins": plugins})

    def _custom_toggle_plugin(self):
        body = self._read_body()
        if isinstance(body, str):
            body = json.loads(body)
        pid = body.get("id", "")
        enabled = body.get("enabled", True)
        state = _get_plugin_enabled_state()
        state[pid] = enabled
        _set_plugin_enabled_state(state)
        self._send_json({"ok": True})

    def _handle_custom_route(self, plugin_id, sub_path, method):
        plugin = _custom_plugins.get(plugin_id)
        if not plugin:
            self._send_json({"ok": False, "error": "plugin not found"}, status=404)
            return

        # Built-in file serving for plugin assets
        if sub_path == "template":
            tpl = plugin["path"] / "template.html"
            if tpl.exists():
                self._send_text(tpl.read_text("utf-8"), content_type="text/html; charset=utf-8")
            else:
                self._send_json({"ok": False, "error": "no template"}, status=404)
            return
        if sub_path == "main.js":
            js = plugin["path"] / "main.js"
            if js.exists():
                self._send_text(js.read_text("utf-8"), content_type="application/javascript; charset=utf-8")
            else:
                self._send_json({"ok": False, "error": "no js"}, status=404)
            return
        if sub_path == "style.css":
            css = plugin["path"] / "style.css"
            if css.exists():
                self._send_text(css.read_text("utf-8"), content_type="text/css; charset=utf-8")
            else:
                self._send_json({"ok": False, "error": "no css"}, status=404)
            return

        # Delegate to routes module
        mod = plugin.get("routes_module")
        if not mod:
            self._send_json({"ok": False, "error": "no routes"}, status=404)
            return

        handler_name = f"handle_{method.lower()}"
        handler = getattr(mod, handler_name, None)
        if not handler:
            self._send_json({"ok": False, "error": f"method {method} not supported"}, status=405)
            return

        try:
            handler(sub_path, self, get_conn)
        except Exception as e:
            logger.error("[Plugin:%s] Error in %s %s: %s", plugin_id, method, sub_path, e)
            self._send_json({"ok": False, "error": str(e)}, status=500)

    def _ai_get_providers(self):
        """Return list of providers with configured API keys."""
        providers = _get_available_providers()
        self._send_json({"ok": True, "providers": providers})

    def _dev_get_ai_keys(self):
        """Return configured AI API keys (masked)."""
        keys = _get_ai_api_keys()
        masked = {}
        for pid, key in keys.items():
            if key and len(key) > 8:
                masked[pid] = key[:4] + "..." + key[-4:]
            elif key:
                masked[pid] = "****"
            else:
                masked[pid] = ""
        self._send_json({"ok": True, "keys": masked, "providers": list(AI_PROVIDERS.keys())})

    def _dev_save_ai_keys(self):
        body = self._read_body()
        data = json.loads(body)
        keys = data.get("keys", {})
        # Merge: only update non-empty values, keep existing for empty ones
        existing = _get_ai_api_keys()
        for pid in AI_PROVIDERS:
            val = keys.get(pid, "").strip()
            if val:
                existing[pid] = val
            # If empty string sent, remove the key
            elif pid in keys and val == "":
                existing.pop(pid, None)
        conn = get_conn()
        conn.execute(
            "INSERT OR REPLACE INTO dev_settings (key, value, encrypted) VALUES ('ai_api_keys', ?, 0)",
            (json.dumps(existing, ensure_ascii=False),),
        )
        conn.commit()
        conn.close()
        logger.info("[AI] API keys updated: %s", [p for p in existing if existing.get(p)])
        self._send_json({"ok": True})

    def _dev_get_modules(self):
        # Reading module config does not require auth
        conn = get_conn()
        row = conn.execute("SELECT value, encrypted FROM dev_settings WHERE key='modules_config'").fetchone()
        conn.close()
        if row:
            value = row[0]
            if row[1] == 1:
                # Encrypted state — decryption needed but key is unknown in session, so use plaintext storage
                pass
            self._send_json({"ok": True, "modules": json.loads(value)})
        else:
            default_modules = {
                "mock": {"log_fetch_limit": 200, "log_max_limit": 1000},
                "csv": {"default_col_width": 120, "min_col_width": 40, "encodings": ["utf-8", "euc-kr", "shift_jis", "iso-8859-1"]},
                "markdown": {"debounce_ms": 300, "min_pane_px": 120, "max_versions": 30},
            }
            self._send_json({"ok": True, "modules": default_modules})

    def _dev_save_modules(self):
        if not self._dev_check_auth():
            return
        body = self._read_body()
        data = json.loads(body)
        modules = data.get("modules", {})
        conn = get_conn()
        conn.execute(
            "INSERT OR REPLACE INTO dev_settings (key, value, encrypted) VALUES ('modules_config', ?, 0)",
            (json.dumps(modules, ensure_ascii=False),),
        )
        conn.commit()
        conn.close()
        self._send_json({"ok": True})

    # ── CDN library sync ──
    def _dev_cdn_status(self):
        """Check local existence and version info for each CDN library"""
        if not self._dev_check_auth():
            return
        meta = _load_vendor_meta()
        items = []
        for lib in CDN_LIBS:
            local_path = _VENDOR_DIR / lib["file"]
            saved = meta.get(lib["name"], {})
            items.append({
                "name": lib["name"],
                "file": lib["file"],
                "currentVersion": saved.get("version", ""),
                "configVersion": lib["version"],
                "exists": local_path.exists(),
                "size": local_path.stat().st_size if local_path.exists() else 0,
            })
        self._send_json({"ok": True, "libs": items})

    def _dev_cdn_check_latest(self):
        """Query latest version of each library from npm registry"""
        if not self._dev_check_auth():
            return
        logger.info("[CDN] Starting latest version check")
        meta = _load_vendor_meta()
        checked = set()
        results = []
        for lib in CDN_LIBS:
            npm_name = lib.get("npm", "")
            if not npm_name or npm_name in checked:
                continue
            checked.add(npm_name)
            latest = _check_npm_latest(npm_name)
            saved_ver = meta.get(lib["name"], {}).get("version", "")
            results.append({
                "npm": npm_name,
                "current": saved_ver or lib["version"],
                "latest": latest or "lookup failed",
            })
        logger.info(f"[CDN] Latest version check complete: {len(results)} packages")
        self._send_json({"ok": True, "results": results})

    def _dev_cdn_sync(self):
        """Download CDN libraries locally (latest or currently configured version)"""
        if not self._dev_check_auth():
            return
        raw = self._read_body()
        try:
            payload = json.loads(raw) if raw.strip() else {}
        except json.JSONDecodeError:
            payload = {}
        use_latest = payload.get("useLatest", False)

        _VENDOR_DIR.mkdir(parents=True, exist_ok=True)
        meta = _load_vendor_meta()
        results = []

        # Query latest versions (if use_latest)
        latest_versions = {}
        if use_latest:
            checked = set()
            for lib in CDN_LIBS:
                npm_name = lib.get("npm", "")
                if not npm_name or npm_name in checked:
                    continue
                checked.add(npm_name)
                ver = _check_npm_latest(npm_name)
                if ver:
                    latest_versions[npm_name] = ver

        for lib in CDN_LIBS:
            local_path = _VENDOR_DIR / lib["file"]
            version = lib["version"]
            if use_latest and lib.get("npm") in latest_versions:
                version = latest_versions[lib["npm"]]
            url = lib["url_tpl"].format(v=version)
            try:
                logger.info(f"[CDN] Downloading: {lib['name']} v{version} ({url})")
                req = urllib.request.Request(url, headers={"User-Agent": "dev-tools/1.0"})
                with urllib.request.urlopen(req, timeout=30) as resp:
                    data = resp.read()
                local_path.write_bytes(data)
                meta[lib["name"]] = {"version": version, "file": lib["file"], "size": len(data)}
                logger.info(f"[CDN] Done: {lib['name']} v{version} ({len(data):,} bytes)")
                results.append({"name": lib["name"], "file": lib["file"], "ok": True, "size": len(data), "version": version})
            except Exception as e:
                logger.error(f"[CDN] Failed: {lib['name']} - {e}")
                results.append({"name": lib["name"], "file": lib["file"], "ok": False, "error": str(e)})

        _save_vendor_meta(meta)
        success = sum(1 for r in results if r["ok"])
        self._send_json({"ok": True, "results": results, "summary": f"{success}/{len(results)} complete"})

    # ══════════════════════════════════
    # Version management API
    # ══════════════════════════════════
    def _dev_get_version(self):
        """Return current version, latest remote version, and update availability."""
        version_file = ROOT / "VERSION"
        current = version_file.read_text().strip() if version_file.exists() else "unknown"

        latest = current
        update_available = False
        try:
            result = subprocess.run(
                ["git", "fetch", "origin", "--tags", "-q"],
                cwd=str(ROOT), capture_output=True, text=True, timeout=30
            )
            result = subprocess.run(
                ["git", "tag", "-l", "v*", "--sort=-version:refname"],
                cwd=str(ROOT), capture_output=True, text=True, timeout=10
            )
            if result.returncode == 0 and result.stdout.strip():
                latest_tag = result.stdout.strip().split("\n")[0]
                latest = latest_tag.lstrip("v")
                if latest != current:
                    update_available = True
        except Exception as e:
            logger.warning(f"[Version] Failed to check remote tags: {e}")

        self._send_json({
            "ok": True,
            "current": current,
            "latest": latest,
            "update_available": update_available,
        })

    def _dev_get_releases(self):
        """Fetch release notes from GitHub API."""
        try:
            result = subprocess.run(
                ["git", "remote", "get-url", "origin"],
                cwd=str(ROOT), capture_output=True, text=True, timeout=5
            )
            remote_url = result.stdout.strip()
            m = re.search(r"github\.com[:/](.+?)(?:\.git)?$", remote_url)
            if not m:
                self._send_json({"ok": False, "error": "not_github"})
                return
            repo_slug = m.group(1)

            api_url = f"https://api.github.com/repos/{repo_slug}/releases?per_page=20"
            req = urllib.request.Request(api_url, headers={
                "Accept": "application/vnd.github+json",
                "User-Agent": "dev-tools/1.0",
            })
            with urllib.request.urlopen(req, timeout=10) as resp:
                releases = json.loads(resp.read().decode())

            items = []
            for r in releases:
                items.append({
                    "tag": r.get("tag_name", ""),
                    "name": r.get("name", ""),
                    "body": r.get("body", ""),
                    "published_at": r.get("published_at", ""),
                    "html_url": r.get("html_url", ""),
                })
            self._send_json({"ok": True, "releases": items})
        except Exception as e:
            logger.warning(f"[Version] Failed to fetch releases: {e}")
            self._send_json({"ok": False, "error": str(e)})

    def _dev_update(self):
        """Pull latest changes from origin main."""
        if not self._dev_check_auth():
            return

        # Check for local changes
        try:
            result = subprocess.run(
                ["git", "status", "--porcelain"],
                cwd=str(ROOT), capture_output=True, text=True, timeout=10
            )
            if result.stdout.strip():
                files = [line.strip() for line in result.stdout.strip().split("\n") if line.strip()]
                self._send_json({"ok": False, "error": "local_changes", "files": files})
                return
        except Exception as e:
            self._send_json({"ok": False, "error": str(e)}, 500)
            return

        # Pull latest
        try:
            result = subprocess.run(
                ["git", "pull", "--ff-only", "origin", "main"],
                cwd=str(ROOT), capture_output=True, text=True, timeout=60
            )
            if result.returncode == 0:
                version_file = ROOT / "VERSION"
                new_version = version_file.read_text().strip() if version_file.exists() else "unknown"
                self._send_json({"ok": True, "version": new_version})
            else:
                # Attempt to abort if merge in progress
                subprocess.run(
                    ["git", "merge", "--abort"],
                    cwd=str(ROOT), capture_output=True, text=True, timeout=10
                )
                self._send_json({"ok": False, "error": "merge_conflict"})
        except Exception as e:
            self._send_json({"ok": False, "error": str(e)}, 500)


if __name__ == "__main__":
    load_env_file(ROOT / ".env")
    parser = argparse.ArgumentParser(description="Test Mock Server")
    parser.add_argument("--host", default=os.getenv("APP_HOST", DEFAULT_HOST))
    parser.add_argument("--port", type=int, default=int(os.getenv("APP_PORT", DEFAULT_PORT)))
    parser.add_argument(
        "--no-open",
        action="store_true",
        help="Disable automatic browser open at startup",
    )
    args = parser.parse_args()

    init_db()
    discover_custom_plugins()
    try:
        server = ThreadingHTTPServer((args.host, args.port), MockHandler)
    except OSError as e:
        if getattr(e, "errno", None) == 48:
            print(f"Port {args.port} is already in use on {args.host}.")
            print(f"Try: lsof -nP -iTCP:{args.port} -sTCP:LISTEN")
            print("Or run with another port, e.g.: python3 server.py --port 9090")
        raise
    ui_url = f"http://{args.host}:{args.port}"
    logger.info(f"Server started: {ui_url}")
    print(f"Mock server running at {ui_url}")
    print("UI: /")
    print("Dynamic mock endpoints: any path except /api/*")
    print(f"Logs: {_LOG_DIR}")
    auto_open = parse_bool(os.getenv("APP_AUTO_OPEN", "true"), default=True) and not args.no_open
    if auto_open:
        try:
            webbrowser.open(ui_url)
        except Exception:
            pass
    try:
        server.serve_forever()
    except KeyboardInterrupt:
        pass
    finally:
        server.server_close()
